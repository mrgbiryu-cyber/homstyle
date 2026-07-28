from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Sequence

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "ppt_assets_combination_20260728"
OUTPUT_PATH = ROOT / "홈스타일_조합상품_3D규격_고객결정안_2026-07-28.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "맑은 고딕"
REPORT_DATE = "2026. 07. 28"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(
        int(value[0:2], 16),
        int(value[2:4], 16),
        int(value[4:6], 16),
    )


WHITE = rgb("#FFFFFF")
BG = rgb("#F5F6F8")
INK = rgb("#202124")
MUTED = rgb("#69727D")
LINE = rgb("#D9DEE5")
LG_RED = rgb("#A50034")
LG_RED_DARK = rgb("#730025")
LG_PINK = rgb("#FBE7EE")
BLUE = rgb("#1A73E8")
BLUE_DARK = rgb("#1557B0")
BLUE_LIGHT = rgb("#E8F0FE")
GREEN = rgb("#188038")
GREEN_LIGHT = rgb("#E6F4EA")
AMBER = rgb("#E37400")
AMBER_LIGHT = rgb("#FEF3E2")
PURPLE = rgb("#6F42C1")
PURPLE_LIGHT = rgb("#F0EAFB")
TEAL = rgb("#0B7B83")
TEAL_LIGHT = rgb("#E3F4F5")
RED = rgb("#C5221F")
RED_LIGHT = rgb("#FCE8E6")
GRAY_DARK = rgb("#3C4043")


@dataclass(frozen=True)
class PatternCase:
    code: str
    count: int
    section: str
    title: str
    product_id: str
    product_name: str
    image_name: str
    question: str
    choices: tuple[str, str, str]
    recommended_choice: str
    recommendation: str
    dimension_rule: str

    @property
    def product_url(self) -> str:
        return f"https://homestyle.lge.co.kr/item?productId={self.product_id}"


CASES: tuple[PatternCase, ...] = (
    PatternCase(
        "N01~N06",
        34,
        "A. 비조합·단일상품",
        "‘+’가 있어도 하나의 제품인가?",
        "G25070001781",
        "LEXON AUDIO 미나 L 무드등 오디오 램프 조명+스피커",
        "G25070001781.jpg",
        "조명+스피커처럼 기능·모델명에 포함된 ‘+’를 조합상품으로 볼 것인가?",
        (
            "A. 비조합 단일 Asset — 제품 전체 W/D/H 1세트",
            "B. 조명·스피커를 각각 component_seq로 분리",
            "C. 단일 Asset과 내부 기능 정보를 함께 보관",
        ),
        "A",
        "‘+’가 기능·모델명·색상·옵션 연결이면 단일상품으로 처리합니다.",
        "숫자를 합산하지 않고 완제품 외곽 W/D/H만 제공합니다.",
    ),
    PatternCase(
        "R01",
        1,
        "A. 비조합·단일상품",
        "내부 수납 구성이 여러 개면 분리할까?",
        "G26050042684",
        "쿠시노 950폭 수납장(PL박스+책장형)",
        "G26050042684.jpg",
        "PL박스와 책장형 수납부처럼 본체 안에 결합된 요소를 별도 상품으로 볼 것인가?",
        (
            "A. 조립된 수납장 1개 — 완제품 W/D/H",
            "B. PL박스·책장부를 별도 component_seq로 분리",
            "C. 완제품 1행 + 내부 구성은 메타정보로 보관",
        ),
        "C",
        "3D Asset은 완제품 1개로 만들고 내부 구성은 설명 정보로 남깁니다.",
        "외곽 W/D/H 1세트만 사용하며 내부 칸 치수는 대표 규격에서 제외합니다.",
    ),
    PatternCase(
        "Y01",
        10,
        "B. 침대·키즈 구성",
        "침대+협탁+침대는 각각 줄까, 한 세트로 줄까?",
        "G25110023368",
        "D2178A 뷰티레스트 자스민 퀸+슈퍼싱글 침대+협탁 101",
        "G25110023368.png",
        "크기가 다른 침대 2개와 협탁을 하나의 W/D/H로 합칠 것인가?",
        (
            "A. 침대 Q / 침대 SS / 협탁을 각각 별도 행",
            "B. 사진 배치 기준으로 조합 전체 외곽 W/D/H 1행",
            "C. 세트 부모 1행 + 구성품별 W/D/H를 모두 제공",
        ),
        "C",
        "배치가 변할 수 있으므로 구성품별 치수를 기본으로 하고 세트 관계를 함께 제공합니다.",
        "기계적 3축 합산 금지. 고정 나란히 배치일 때만 W=합, D/H=최대값입니다.",
    ),
    PatternCase(
        "Y02",
        17,
        "B. 침대·키즈 구성",
        "침대와 협탁은 독립 Asset인가?",
        "G25110023359",
        "D2178A 뷰티레스트 자스민 킹오브킹 침대+협탁 101 2개",
        "G25110023359.png",
        "침대 옆에서 자유롭게 이동하는 협탁을 침대 전체 규격에 포함할 것인가?",
        (
            "A. 침대 1행 + 협탁 1행(수량은 별도)",
            "B. 침대와 협탁을 합친 전체 외곽 W/D/H",
            "C. 침대 대표 규격만 제공하고 협탁은 제외",
        ),
        "A",
        "독립 이동 가능한 협탁은 별도 component_seq로 제공합니다.",
        "침대와 협탁 치수를 각각 보존하며 침대 폭에 협탁 폭을 더하지 않습니다.",
    ),
    PatternCase(
        "Y03",
        4,
        "B. 침대·키즈 구성",
        "침대 패널·가드는 본체 치수에 포함할까?",
        "G26020033567",
        "베른 가죽침대 Q + 패널 60cm 2EA",
        "G26020033567.jpg",
        "침대에 고정되는 패널·가드를 별도 Asset으로 만들지, 조립된 외곽에 포함할지 결정이 필요합니다.",
        (
            "A. 침대와 패널·가드를 각각 별도 행",
            "B. 설치 완료 상태의 침대 전체 W/D/H 1행",
            "C. 전체 W/D/H + 패널·가드 부속품 메타정보",
        ),
        "C",
        "고정 설치되는 구조물은 조립 후 전체 외곽에 포함하고 부속 관계를 남깁니다.",
        "전체 폭은 설치 방향·개수에 따라 계산하며 패널 폭을 무조건 더하지 않습니다.",
    ),
    PatternCase(
        "Y05",
        14,
        "B. 침대·키즈 구성",
        "유아동 침대+매트리스는 한 Asset인가?",
        "G26030036471",
        "한샘 티오 그로우 슬라이딩 딥 수납침대 SS+샘키즈 매트",
        "G26030036471_alt3.jpg",
        "프레임과 매트리스를 3D에서 하나로 만들지, 각각 교체 가능한 구성품으로 만들지 결정이 필요합니다.",
        (
            "A. 프레임과 매트리스를 각각 별도 행",
            "B. 침대 시스템 1개 — 설치 상태 전체 W/D/H",
            "C. 침대 1행 + 매트리스 규격을 부속 메타정보로 제공",
        ),
        "C",
        "기존 침대 정책과 맞춰 단일 Asset으로 처리하되 매트리스 규격은 추적 가능하게 보존합니다.",
        "대표 W/D/H는 침대 외곽, 매트리스 W/D/H는 구성 메타정보로 구분합니다.",
    ),
    PatternCase(
        "Y06",
        1,
        "B. 침대·키즈 구성",
        "유아동 침대+가드(+책상)는 혼합 처리할까?",
        "G26030036470",
        "한샘 티오 그로우 슬라이딩 딥 수납침대 SS+가드",
        "G26030036470_alt3.jpg",
        "본체에 고정되는 가드와 독립 이동 가능한 책상·수납장을 같은 방식으로 처리할 수 없습니다.",
        (
            "A. 모든 구성품을 각각 별도 Asset",
            "B. 모든 구성품을 하나의 조합 전체 Asset",
            "C. 침대+고정 가드는 통합, 이동 가구는 별도 행",
        ),
        "C",
        "고정 여부에 따라 혼합 처리하는 방식이 실제 배치와 가장 잘 맞습니다.",
        "통합부는 조립 외곽 W/D/H, 이동 가구는 각자 W/D/H를 제공합니다.",
    ),
    PatternCase(
        "Y07",
        12,
        "C. 가구 결합·모듈",
        "식탁+의자·벤치는 조합 전체 치수가 필요한가?",
        "G25080006948",
        "Table + Bench 003",
        "G25080006948.jpg",
        "테이블과 벤치·의자는 사용 중 위치가 바뀌므로 하나의 고정 외곽 규격을 정의하기 어렵습니다.",
        (
            "A. 테이블 / 벤치 / 의자를 각각 별도 행",
            "B. 대표 연출 배치의 전체 외곽 W/D/H",
            "C. 테이블만 대표 규격으로 제공하고 좌석은 수량 정보",
        ),
        "A",
        "구성품별 Asset과 수량을 제공하고 세트 ID로 연결합니다.",
        "테이블·벤치·의자의 W/D/H를 각각 보존하며 전체 합산값은 만들지 않습니다.",
    ),
    PatternCase(
        "Y08",
        1,
        "C. 가구 결합·모듈",
        "상판+다리는 완제품 한 개로 볼까?",
        "G25110022880",
        "베가 상판+위시본 다리 Signature",
        "G25110022880.jpg",
        "조립 후 독립적으로 움직이지 않는 상판과 다리를 별도 Asset으로 전달할 필요가 있는지 결정합니다.",
        (
            "A. 상판과 다리를 각각 별도 행",
            "B. 조립 완료 테이블 1개 — 전체 W/D/H",
            "C. 완제품 1행 + 상판·다리 부품 메타정보",
        ),
        "C",
        "3D 전달은 완제품 1개, 부품 추적이 필요하면 내부 메타정보로 보관합니다.",
        "완제품 외곽 W/D/H를 사용하고 부품 치수를 더해 대표값을 만들지 않습니다.",
    ),
    PatternCase(
        "Y09",
        20,
        "C. 가구 결합·모듈",
        "모듈 소파는 모듈별인가, 조합 완료형인가?",
        "G25120026411",
        "톤카 소파 A1+A0+A1",
        "G25120026411.jpg",
        "판매 SKU가 특정 모듈 조합을 의미할 때 조합 완료형과 개별 모듈 중 어느 단위로 3D를 만들지 결정합니다.",
        (
            "A. A1 / A0 / A1 모듈을 각각 별도 행",
            "B. 해당 SKU의 조합 완료 소파 1개",
            "C. 조합 완료 부모 1행 + 모듈별 치수·순서",
        ),
        "C",
        "렌더링은 SKU별 조합 완료형을 사용하고 재조합을 위해 모듈 순서를 함께 보존합니다.",
        "대표 W/D/H는 조합 완료 외곽값이며 모듈 폭 단순합은 실제 결합값과 대조합니다.",
    ),
    PatternCase(
        "Y10",
        5,
        "D. 부속품·수량",
        "소파+헤드레스트·쿠션은 어디까지 모델링할까?",
        "G25110024764",
        "밴쿠버 패브릭 소파 3.5인 + 헤드레스트 2EA",
        "G25110024764.jpg",
        "탈착 가능한 쿠션·헤드레스트를 독립 3D Asset으로 만들지 고객 정책이 필요합니다.",
        (
            "A. 소파와 부속품을 각각 별도 Asset",
            "B. 부속품 장착 상태의 소파 1개만 제공",
            "C. 소파 본체 1행 + 부속품 종류·수량 메타정보",
        ),
        "C",
        "대표 규격은 소파 본체 기준으로 두고 부속품 종류와 수량을 별도 보존합니다.",
        "헤드레스트 장착 시 높이가 달라지면 option_seq로 장착 상태 규격을 추가합니다.",
    ),
    PatternCase(
        "Y11",
        10,
        "C. 가구 결합·모듈",
        "수납장 80+120cm는 전체 200cm인가?",
        "G26030035534",
        "한샘 클린트 모던 높은 거실장 200cm(80cm+120cm)",
        "G26030035534.jpg",
        "나란히 고정 설치하는 모듈인지, 각각 독립 배치 가능한 장인지에 따라 규격 단위가 달라집니다.",
        (
            "A. 80cm장 / 120cm장을 각각 별도 행",
            "B. 설치 완료 200cm 수납장 전체 W/D/H",
            "C. 전체 200cm 부모 1행 + 모듈별 W/D/H",
        ),
        "C",
        "상품명이 완성 폭을 명시하면 전체 규격을 제공하고 모듈별 치수도 추적합니다.",
        "나란히 밀착 설치 시 W=80+120cm, D/H는 모듈 최대값입니다.",
    ),
    PatternCase(
        "Y12",
        4,
        "D. 부속품·수량",
        "행거+후크·옷걸이는 별도 치수가 필요한가?",
        "G26070047242",
        "ON 폴 이동식 바지걸이 1단 행거 L + 후크20P",
        "G26070047242.jpg",
        "작은 후크·옷걸이를 독립 3D Asset으로 만들지, 수량 메타정보로만 제공할지 결정합니다.",
        (
            "A. 행거와 후크를 각각 별도 Asset",
            "B. 행거 본체 W/D/H + 후크 수량만 제공",
            "C. 후크 장착 상태 전체 외곽 W/D/H",
        ),
        "B",
        "소형 부속품은 3D 범위에서 제외하고 본체 규격과 수량 메타만 제공합니다.",
        "후크 개수로 행거 W/D/H를 곱하거나 합산하지 않습니다.",
    ),
    PatternCase(
        "Y13",
        12,
        "C. 가구 결합·모듈",
        "옷장 모듈은 전체 폭과 개별 폭을 모두 줄까?",
        "G25070005769",
        "일상고백 클라로 옷장 280cm(이불장+거울장+행거장+서랍장)",
        "G25070005769.jpg",
        "설치 후 연속된 옷장으로 쓰지만 모듈 구성이 옵션별로 달라 전체·개별 규격을 함께 관리할지 결정합니다.",
        (
            "A. 모듈별 W/D/H만 제공",
            "B. 설치 완료 옷장 전체 W/D/H만 제공",
            "C. 전체 부모 1행 + 모듈별 W/D/H와 순서",
        ),
        "C",
        "전체 설치 폭과 모듈 구성을 모두 제공해야 옵션 변경과 3D 조립을 함께 지원할 수 있습니다.",
        "전체 W는 설치 순서의 폭 합, D/H는 최대값이며 틈·몰딩은 명시값만 포함합니다.",
    ),
    PatternCase(
        "Y14",
        3,
        "D. 부속품·수량",
        "1+1 상품의 폭을 두 배로 줄까?",
        "G26040039060",
        "한샘 PP 스퀘어 스툴 1+1",
        "G26040039060.jpg",
        "동일 스툴 2개는 배치가 정해져 있지 않으므로 규격과 수량을 분리해야 합니다.",
        (
            "A. 스툴 W/D/H 1세트 + quantity=2",
            "B. 두 개를 나란히 놓은 W×2 전체 규격",
            "C. 같은 치수의 구성품 행을 2개 생성",
        ),
        "A",
        "동일 Asset 하나와 수량을 전달하는 방식이 중복 제작을 막습니다.",
        "W/D/H는 1개 기준이며 수량만 2로 기록합니다.",
    ),
    PatternCase(
        "Y15",
        1,
        "D. 부속품·수량",
        "책상 액세서리 팩은 무엇을 3D로 만들까?",
        "G26030036463",
        "한샘 각도조절 FRIENDS PACK(콘센트·조명·데스크패드+자석바)",
        "G26030036463_alt2.jpg",
        "상품 자체가 여러 액세서리의 묶음이므로 패키지 전체 W/D/H는 실사용 형상을 설명하지 못합니다.",
        (
            "A. 모델링 대상 부속품을 각각 별도 행",
            "B. 패키지 박스 전체 W/D/H 1행",
            "C. 3D 대상에서 제외하고 호환 정보만 제공",
        ),
        "A",
        "3D로 표현할 부속품만 component_seq로 분리하고 포장 규격은 제외합니다.",
        "콘센트·조명·패드 치수를 각각 사용하며 서로 합산하지 않습니다.",
    ),
)


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor, width: float = 0.8) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 12,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
    linespacing: float | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    if linespacing is not None:
        p.line_spacing = linespacing
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return shape


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    radius: bool = True,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    set_line(shape, line)
    return shape


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    color: RGBColor,
    size: float = 9,
):
    add_box(slide, x, y, w, 0.34, fill=fill, line=fill)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_slide_base(
    prs: Presentation,
    section: str,
    title: str,
    page: int,
    subtitle: str = "",
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_text(slide, section.upper(), 0.68, 0.26, 6.3, 0.25, size=8.8, color=LG_RED, bold=True)
    add_text(slide, title, 0.68, 0.59, 11.9, 0.55, size=23, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.70, 1.13, 11.8, 0.32, size=9.5, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.47), Inches(11.95), Inches(0.012)
    )
    set_fill(line, LINE)
    line.line.fill.background()
    add_text(slide, f"HomeStyle 3D combination decision · {REPORT_DATE}", 0.7, 7.12, 4.5, 0.16, size=7.2, color=MUTED)
    add_text(slide, f"{page:02d}", 12.35, 7.12, 0.28, 0.16, size=7.2, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def add_picture_fit(slide, image_path: Path, x: float, y: float, w: float, h: float):
    with Image.open(image_path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(x + (w - pw) / 2),
        Inches(y + (h - ph) / 2),
        Inches(pw),
        Inches(ph),
    )
    return picture


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.end_arrowhead = True
    return line


def make_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(SLIDE_H)
    )
    set_fill(accent, LG_RED)
    accent.line.fill.background()
    add_pill(slide, "CUSTOMER DECISION DECK · 149 PRODUCTS", 0.86, 0.66, 3.15, fill=LG_RED_DARK, color=WHITE)
    add_text(
        slide,
        "홈스타일 조합상품\n3D 규격 전달 기준 결정안",
        0.86,
        1.31,
        7.4,
        1.52,
        size=29,
        color=WHITE,
        bold=True,
    )
    add_text(
        slide,
        "구성품별 규격 · 조합 전체 규격 · 수량·부속품 처리\n16개 대표 패턴과 실제 상품 이미지로 결정",
        0.89,
        3.22,
        7.15,
        0.92,
        size=14,
        color=rgb("#D9DEE5"),
    )
    add_text(slide, REPORT_DATE, 0.89, 6.83, 1.4, 0.2, size=8.5, color=rgb("#AEB5BD"))
    add_text(slide, "LG HomeStyle Data PoC", 2.35, 6.83, 2.5, 0.2, size=8.5, color=rgb("#AEB5BD"))
    add_text(
        slide,
        "기준 문서 · 조합상품_3D변환_패턴목록_2026-07-27.html",
        0.89,
        6.46,
        6.2,
        0.20,
        size=7.8,
        color=rgb("#AEB5BD"),
    )

    card = add_box(slide, 8.55, 0.72, 4.1, 5.92, fill=rgb("#2B2D31"), line=rgb("#41454B"))
    card.shadow.inherit = False
    samples = (
        ASSET_DIR / "G25110023368.png",
        ASSET_DIR / "G25080006948.jpg",
        ASSET_DIR / "G25120026411.jpg",
    )
    for idx, image_path in enumerate(samples):
        yy = 1.02 + idx * 1.65
        add_box(slide, 8.88, yy, 1.68, 1.34, fill=WHITE, line=WHITE)
        pic = add_picture_fit(slide, image_path, 8.95, yy + 0.07, 1.54, 1.20)
        add_text(
            slide,
            ("침대+협탁+침대", "테이블+벤치", "모듈 소파")[idx],
            10.78,
            yy + 0.18,
            1.42,
            0.35,
            size=10.2,
            color=WHITE,
            bold=True,
        )
        add_text(
            slide,
            ("개별 vs 전체", "분리 Asset", "조합 완료형")[idx],
            10.78,
            yy + 0.65,
            1.42,
            0.25,
            size=8.5,
            color=rgb("#B9C0C8"),
        )
        pic.click_action.hyperlink.address = (
            "https://homestyle.lge.co.kr/item?productId="
            + ("G25110023368", "G25080006948", "G25120026411")[idx]
        )
    add_pill(slide, "16 PATTERNS", 9.0, 6.03, 1.3, fill=LG_RED_DARK, color=WHITE)
    add_text(slide, "고객 승인 후 149개에 일괄 적용", 10.48, 6.07, 1.8, 0.25, size=8.7, color=rgb("#D9DEE5"))


def make_toc(prs: Presentation, page: int) -> None:
    slide = add_slide_base(
        prs,
        "CONTENTS",
        "목차",
        page,
        "149개 상품을 성격이 비슷한 4개 결정군과 16개 대표 패턴으로 정리했습니다.",
    )
    sections = (
        ("A", "비조합·단일상품", "35개", "N01~N06 · R01", LG_RED, LG_PINK),
        ("B", "침대·키즈 구성", "46개", "Y01 · Y02 · Y03 · Y05 · Y06", BLUE, BLUE_LIGHT),
        ("C", "가구 결합·모듈", "55개", "Y07 · Y08 · Y09 · Y11 · Y13", TEAL, TEAL_LIGHT),
        ("D", "부속품·수량", "13개", "Y10 · Y12 · Y14 · Y15", PURPLE, PURPLE_LIGHT),
    )
    positions = ((0.72, 1.78), (6.84, 1.78), (0.72, 4.23), (6.84, 4.23))
    for (label, title, count, codes, accent, light), (x, y) in zip(sections, positions):
        add_box(slide, x, y, 5.78, 1.92, fill=WHITE, line=LINE)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.28), Inches(y + 0.27), Inches(0.58), Inches(0.58)
        )
        set_fill(circle, accent)
        circle.line.fill.background()
        add_text(slide, label, x + 0.28, y + 0.29, 0.58, 0.5, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, x + 1.05, y + 0.23, 3.55, 0.38, size=16, bold=True)
        add_pill(slide, count, x + 4.57, y + 0.24, 0.85, fill=light, color=accent, size=10)
        add_text(slide, codes, x + 1.05, y + 0.85, 4.1, 0.34, size=10, color=accent, bold=True)
        descriptions = (
            "‘+’ 표기가 구성품이 아닌 경우 · 내부 결합 요소",
            "침대·협탁·패널·매트리스·가드 관계",
            "테이블·소파·수납장·옷장 모듈",
            "쿠션·후크·1+1·액세서리 팩",
        )
        add_text(slide, descriptions[sections.index((label, title, count, codes, accent, light))], x + 1.05, y + 1.33, 4.2, 0.28, size=9, color=MUTED)
    add_text(slide, "참고", 0.76, 6.53, 0.55, 0.22, size=8.5, color=GREEN, bold=True)
    add_text(slide, "침대 프레임+매트리스 121개(Y04)는 기존 정책상 단일 Asset으로 이미 확정되어 본 결정 대상에서 제외했습니다.", 1.35, 6.50, 10.65, 0.28, size=9.2, color=MUTED)


def make_rules(prs: Presentation, page: int) -> None:
    slide = add_slide_base(
        prs,
        "DECISION GUIDE",
        "규격은 ‘숫자 합산’이 아니라 배치·고정 방식으로 결정합니다",
        page,
        "아래 네 가지 원칙을 먼저 승인하면 16개 패턴의 판단이 단순해집니다.",
    )
    cards = (
        ("1", "독립 이동", "침대+협탁, 식탁+의자", "구성품별 W/D/H\n전체 합산하지 않음", BLUE, BLUE_LIGHT),
        ("2", "고정 조립", "상판+다리, 붙박이 모듈", "조립 후 외곽 W/D/H\n부품은 메타정보", TEAL, TEAL_LIGHT),
        ("3", "동일 수량", "스툴 1+1, 의자 4P", "W/D/H 1세트\n+ quantity", PURPLE, PURPLE_LIGHT),
        ("4", "옵션 변화", "SS/Q/K, 좌형/우형", "외형이 달라질 때만\noption_seq 분리", AMBER, AMBER_LIGHT),
    )
    xs = (0.72, 3.83, 6.94, 10.05)
    for (num, title, example, output, accent, light), x in zip(cards, xs):
        add_box(slide, x, 1.78, 2.62, 2.24, fill=WHITE, line=LINE)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.20), Inches(1.98), Inches(0.45), Inches(0.45)
        )
        set_fill(circle, accent)
        circle.line.fill.background()
        add_text(slide, num, x + 0.20, 2.0, 0.45, 0.39, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, title, x + 0.82, 1.96, 1.45, 0.35, size=14, bold=True)
        add_text(slide, example, x + 0.20, 2.63, 2.18, 0.28, size=8.6, color=MUTED)
        add_box(slide, x + 0.18, 3.08, 2.26, 0.68, fill=light, line=light)
        add_text(slide, output, x + 0.28, 3.16, 2.02, 0.5, size=9.4, color=accent, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)

    add_box(slide, 0.72, 4.38, 12.0, 1.82, fill=WHITE, line=LINE)
    add_text(slide, "나란히 고정 배치하는 경우의 전체 외곽 계산", 0.98, 4.62, 4.1, 0.31, size=13, bold=True)
    for idx, (x, label, color) in enumerate(((1.05, "W₁", BLUE), (2.07, "W₂", TEAL), (3.09, "W₃", PURPLE))):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(5.14), Inches(0.88), Inches(0.58 + idx * 0.08)
        )
        set_fill(shape, color)
        shape.line.fill.background()
        add_text(slide, label, x, 5.28, 0.88, 0.24, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 4.30, 5.48, 5.05, 5.48, LG_RED)
    add_text(slide, "W전체 = W₁+W₂+W₃ (+명시된 틈)", 5.30, 4.92, 3.1, 0.34, size=12.3, color=LG_RED, bold=True)
    add_text(slide, "D전체 = max(D₁,D₂,D₃)  ·  H전체 = max(H₁,H₂,H₃)", 5.30, 5.42, 4.0, 0.33, size=11.2, color=GRAY_DARK, bold=True)
    add_box(slide, 9.55, 4.72, 2.76, 1.12, fill=RED_LIGHT, line=RED_LIGHT)
    add_text(slide, "주의", 9.78, 4.89, 0.45, 0.24, size=9, color=RED, bold=True)
    add_text(slide, "W·D·H를 모두 더하거나\n사진 속 임의 간격을 포함하지 않습니다.", 10.28, 4.84, 1.74, 0.62, size=9.1, color=RED, bold=True)


def make_overview(prs: Presentation, page: int) -> None:
    slide = add_slide_base(
        prs,
        "DECISION SCOPE",
        "현재 고객 결정 대상은 149개 · 16개 대표 패턴입니다",
        page,
        "패턴 승인을 받으면 같은 구조의 상품에 동일 규칙을 적용할 수 있습니다.",
    )
    data = (
        ("비조합·단일", 35, "완제품 W/D/H", LG_RED, LG_PINK),
        ("침대·키즈", 46, "분리·통합 혼합", BLUE, BLUE_LIGHT),
        ("가구 결합·모듈", 55, "설치 방식 확인", TEAL, TEAL_LIGHT),
        ("부속품·수량", 13, "메타·quantity", PURPLE, PURPLE_LIGHT),
    )
    max_count = max(v for _, v, _, _, _ in data)
    for idx, (label, count, note, accent, light) in enumerate(data):
        y = 1.82 + idx * 1.0
        add_text(slide, label, 0.82, y + 0.13, 1.55, 0.28, size=10.8, bold=True)
        add_box(slide, 2.45, y, 6.65, 0.55, fill=rgb("#E8EAED"), line=rgb("#E8EAED"))
        bar_w = 6.65 * count / max_count
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.45), Inches(y), Inches(bar_w), Inches(0.55)
        )
        set_fill(bar, accent)
        bar.line.fill.background()
        add_text(slide, f"{count}개", 2.62, y + 0.10, 0.72, 0.27, size=11, color=WHITE, bold=True)
        add_pill(slide, note, 9.36, y + 0.10, 2.0, fill=light, color=accent, size=8.5)

    add_box(slide, 0.78, 6.02, 11.85, 0.64, fill=GREEN_LIGHT, line=GREEN_LIGHT)
    add_text(slide, "확정 완료", 1.02, 6.18, 0.8, 0.25, size=9.5, color=GREEN, bold=True)
    add_text(slide, "Y04 침대 프레임+매트리스 121개 → 단일 Asset 정책 유지", 1.90, 6.16, 4.8, 0.26, size=10.2, color=GREEN, bold=True)
    add_text(slide, "이번 회신 대상에는 포함하지 않음", 9.68, 6.18, 2.2, 0.24, size=8.8, color=GREEN)


def make_case_slide(prs: Presentation, case: PatternCase, page: int) -> None:
    slide = add_slide_base(
        prs,
        case.section,
        f"{case.code}  {case.title}",
        page,
        f"해당 패턴 {case.count}개 · 대표 사례 {case.product_id}",
    )
    add_box(slide, 0.70, 1.71, 4.87, 4.92, fill=WHITE, line=LINE)
    add_pill(slide, "실제 상품 이미지", 0.94, 1.93, 1.22, fill=LG_PINK, color=LG_RED, size=8.2)
    image_path = ASSET_DIR / case.image_name
    picture = add_picture_fit(slide, image_path, 0.95, 2.30, 4.36, 3.34)
    picture.click_action.hyperlink.address = case.product_url
    add_text(slide, case.product_id, 0.96, 5.79, 1.32, 0.22, size=8.5, color=BLUE, bold=True)
    product_shape = add_text(slide, case.product_name, 0.96, 6.05, 4.18, 0.40, size=9.1, color=INK, bold=True)
    product_shape.click_action.hyperlink.address = case.product_url
    add_text(slide, "이미지·상품명: HomeStyle 상품 API · 클릭 시 PDP 이동", 0.96, 6.48, 4.15, 0.18, size=6.8, color=MUTED)

    add_box(slide, 5.84, 1.71, 6.79, 0.86, fill=LG_PINK, line=LG_PINK)
    add_text(slide, "고객 결정 질문", 6.10, 1.89, 1.02, 0.24, size=9, color=LG_RED, bold=True)
    add_text(slide, case.question, 7.15, 1.83, 5.13, 0.48, size=10.2, color=INK, bold=True, valign=MSO_ANCHOR.MIDDLE)

    choice_colors = ((BLUE, BLUE_LIGHT), (TEAL, TEAL_LIGHT), (PURPLE, PURPLE_LIGHT))
    for idx, (choice, (accent, light)) in enumerate(zip(case.choices, choice_colors)):
        y = 2.78 + idx * 0.82
        selected = choice.startswith(case.recommended_choice + ".")
        fill = light if selected else WHITE
        line = accent if selected else LINE
        add_box(slide, 5.84, y, 6.79, 0.68, fill=fill, line=line)
        square = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.08), Inches(y + 0.15), Inches(0.34), Inches(0.34)
        )
        set_fill(square, accent if selected else WHITE)
        set_line(square, accent, 1.2)
        if selected:
            add_text(slide, "✓", 6.08, y + 0.14, 0.34, 0.29, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, choice, 6.58, y + 0.12, 5.42, 0.42, size=9.7, color=INK, bold=selected, valign=MSO_ANCHOR.MIDDLE)
        if selected:
            add_pill(slide, "기본 제안", 11.62, y + 0.17, 0.76, fill=accent, color=WHITE, size=7.2)

    add_box(slide, 5.84, 5.42, 6.79, 1.06, fill=WHITE, line=LINE)
    add_text(slide, "권고안", 6.10, 5.61, 0.54, 0.24, size=8.8, color=GREEN, bold=True)
    add_text(slide, case.recommendation, 6.76, 5.55, 5.50, 0.38, size=9.6, color=INK, bold=True)
    add_text(slide, "규격 산출", 6.10, 6.06, 0.62, 0.23, size=8.5, color=AMBER, bold=True)
    add_text(slide, case.dimension_rule, 6.76, 6.00, 5.48, 0.38, size=8.9, color=MUTED)
    add_text(slide, "고객 선택", 5.96, 6.68, 0.70, 0.20, size=8.5, color=LG_RED, bold=True)
    add_text(slide, "□ A     □ B     □ C     □ 기타: ____________________", 6.76, 6.64, 5.45, 0.26, size=9.4, color=GRAY_DARK)


def make_response_slide(prs: Presentation, page: int) -> None:
    slide = add_slide_base(
        prs,
        "RESPONSE",
        "회신 시 아래 6개 원칙과 예외 패턴만 확정해 주세요",
        page,
        "각 대표 사례 장표의 A/B/C 선택을 회신하면 149개 상품에 규칙을 일괄 적용할 수 있습니다.",
    )
    questions = (
        ("01", "이동 가능한 구성품", "구성품별 W/D/H + set_id", BLUE, BLUE_LIGHT),
        ("02", "고정 조립 구성품", "조립 완료 외곽 W/D/H", TEAL, TEAL_LIGHT),
        ("03", "부모+구성품 동시 제공", "필요 패턴만 parent/component 병행", PURPLE, PURPLE_LIGHT),
        ("04", "소형 부속품 모델링 범위", "쿠션·후크·콘센트 포함 여부", AMBER, AMBER_LIGHT),
        ("05", "수량 상품", "규격 1세트 + quantity", GREEN, GREEN_LIGHT),
        ("06", "옵션 규격 분리", "외형이 달라질 때만 option_seq", LG_RED, LG_PINK),
    )
    for idx, (num, title, body, accent, light) in enumerate(questions):
        col = idx % 2
        row = idx // 2
        x = 0.74 + col * 6.08
        y = 1.80 + row * 1.47
        add_box(slide, x, y, 5.70, 1.15, fill=WHITE, line=LINE)
        add_pill(slide, num, x + 0.22, y + 0.18, 0.48, fill=accent, color=WHITE, size=8.2)
        add_text(slide, title, x + 0.88, y + 0.16, 2.2, 0.30, size=12, bold=True)
        add_text(slide, body, x + 0.88, y + 0.57, 3.82, 0.28, size=9.2, color=MUTED)
        add_text(slide, "□ 승인   □ 수정", x + 4.22, y + 0.19, 1.15, 0.26, size=8.5, color=accent, bold=True)
    add_box(slide, 0.74, 6.25, 11.78, 0.52, fill=LG_PINK, line=LG_PINK)
    add_text(slide, "회신 형식 예시", 0.98, 6.40, 1.05, 0.22, size=8.7, color=LG_RED, bold=True)
    add_text(slide, "Y01=C, Y02=A, Y03=C … / 예외 상품은 product_id와 희망 규격 단위를 함께 기재", 2.16, 6.36, 8.95, 0.28, size=9.6, color=INK, bold=True)


def build() -> Path:
    assert sum(case.count for case in CASES) == 149
    missing = [case.image_name for case in CASES if not (ASSET_DIR / case.image_name).exists()]
    cover_missing = [
        name
        for name in ("G25110023368.png", "G25080006948.jpg", "G25120026411.jpg")
        if not (ASSET_DIR / name).exists()
    ]
    if missing or cover_missing:
        raise FileNotFoundError(
            "대표 이미지가 없습니다: " + ", ".join(sorted(set(missing + cover_missing)))
        )

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "홈스타일 조합상품 3D 규격 고객 결정안"
    prs.core_properties.subject = "149개 조합상품의 구성품·전체규격·수량·부속품 처리 기준"
    prs.core_properties.author = "LG HomeStyle Data PoC"
    prs.core_properties.keywords = "HomeStyle, 3D, 조합상품, WDH, 고객결정"

    make_cover(prs)
    page = 2
    make_toc(prs, page)
    page += 1
    make_rules(prs, page)
    page += 1
    make_overview(prs, page)
    page += 1
    for case in CASES:
        make_case_slide(prs, case, page)
        page += 1
    make_response_slide(prs, page)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    output = build()
    print(output)
