from __future__ import annotations

import sqlite3
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DB_PATH = ROOT / "homestyle_bulk_run" / "homestyle_bulk.sqlite"
OUTPUT_PATH = ROOT / "홈스타일_상품정보_OCR_개발인수인계_보고자료_2026-07-27.pptx"
WORKBOOK_PATH = ROOT / "홈스타일_비음영대상군_전체상품_요구필드_대량결과_패턴상태.xlsx"
REPORT_DATE = "2026.07.27"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "맑은 고딕"


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


WHITE = rgb("#FFFFFF")
BG = rgb("#F6F7F9")
INK = rgb("#202124")
MUTED = rgb("#69727D")
LINE = rgb("#D9DEE5")
LG_RED = rgb("#A50034")
LG_RED_DARK = rgb("#6F0023")
LG_PINK = rgb("#FBE7EE")
BLUE = rgb("#1A73E8")
BLUE_DARK = rgb("#1557B0")
BLUE_LIGHT = rgb("#E8F0FE")
GREEN = rgb("#188038")
GREEN_LIGHT = rgb("#E6F4EA")
AMBER = rgb("#E37400")
AMBER_LIGHT = rgb("#FEF3E2")
RED = rgb("#C5221F")
RED_LIGHT = rgb("#FCE8E6")
GRAY_LIGHT = rgb("#ECEFF3")
TEAL = rgb("#0B7B83")
TEAL_LIGHT = rgb("#E3F4F5")
PURPLE = rgb("#7E57C2")
PURPLE_LIGHT = rgb("#F0EAFB")


@dataclass
class Stats:
    total: int
    locked: int
    source_confirmed: int
    rule_resolved: int
    comparison: int
    no_candidate: int
    mandatory_pass: int
    reinforcement: int
    color_missing: int
    size_missing: int
    categories: int
    sources: int
    pass1_products: int
    pass1_images: int
    pass2_products: int
    pass2_observations: int
    targeted_crop_total: int
    comparison_candidates: int
    combination_title_candidates: int
    combination_high_confidence: int
    sofa_stool_combinations: int
    component_dimensions_fully_confirmed: int
    component_rows: int
    bed_single_asset: int
    workbook_columns: int
    workbook_sheets: int
    pattern_no_check: int
    pattern_dimension: int
    pattern_combination: int
    pattern_both: int
    dimension_patterns: dict[str, int]
    combination_pattern_groups: dict[str, int]
    remaining: list[tuple[str, str, int]]

    @property
    def analyzable(self) -> int:
        return self.locked + self.comparison

    @property
    def analyzable_pct(self) -> float:
        return self.analyzable / self.total * 100

    @property
    def pass_pct(self) -> float:
        return self.mandatory_pass / self.total * 100

    @property
    def pattern_check_needed(self) -> int:
        return self.pattern_dimension + self.pattern_combination + self.pattern_both

    @property
    def pattern_check_pct(self) -> float:
        return self.pattern_check_needed / self.total * 100


def load_stats() -> Stats:
    # Excel과 PPT의 상태 수치가 갈라지지 않도록 동일한 workbook 산출 로직을
    # read-only로 호출해 최신 상태 메타데이터를 가져온다.
    from build_homestyle_bulk_workbook import build_rows as build_workbook_rows

    _, workbook_meta = build_workbook_rows()
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    progress = conn.execute("SELECT * FROM vw_dimension_progress_authoritative").fetchone()
    mandatory = {
        row["final_status"]: row["product_count"]
        for row in conn.execute(
            "SELECT final_status, product_count FROM vw_mandatory_pass_current_summary"
        )
    }
    current_snapshot = conn.execute(
        "SELECT snapshot_id FROM stg_mandatory_pass WHERE is_current=1 LIMIT 1"
    ).fetchone()[0]
    missing = {}
    for column in ("color_ok", "size_wdh_ok"):
        missing[column] = conn.execute(
            f"""
            SELECT SUM(CASE WHEN {column}=0 THEN 1 ELSE 0 END)
            FROM stg_mandatory_pass
            WHERE snapshot_id=?
            """,
            (current_snapshot,),
        ).fetchone()[0]

    remaining = [
        (row["reason_code"], row["reason_name"], row["product_count"])
        for row in conn.execute(
            """
            SELECT reason_code, reason_name, product_count
            FROM vw_dimension_remaining_reason_summary
            ORDER BY product_count DESC
            """
        )
    ]

    def scalar(query: str) -> int:
        return int(conn.execute(query).fetchone()[0])

    stats = Stats(
        total=int(progress["total_products"]),
        locked=int(progress["locked_resolved"]),
        source_confirmed=int(progress["source_confirmed"]),
        rule_resolved=int(progress["rule_resolved"]),
        comparison=int(progress["comparison_provided"]),
        no_candidate=int(progress["no_candidate"]),
        mandatory_pass=int(mandatory["PASS"]),
        reinforcement=int(mandatory["보강대상"]),
        color_missing=missing["color_ok"],
        size_missing=missing["size_wdh_ok"],
        categories=scalar("SELECT COUNT(*) FROM categories"),
        sources=scalar("SELECT COUNT(*) FROM sources"),
        pass1_products=scalar(
            "SELECT COUNT(DISTINCT product_id) FROM stg_dimension_scan_pass1_product_image"
        ),
        pass1_images=scalar("SELECT COUNT(*) FROM stg_dimension_scan_pass1_product_image"),
        pass2_products=scalar("SELECT COUNT(*) FROM stg_dimension_scan_pass2_product"),
        pass2_observations=scalar("SELECT COUNT(*) FROM stg_dimension_scan_pass2_observation"),
        targeted_crop_total=scalar("SELECT COUNT(*) FROM stg_dimension_targeted_ocr_crop"),
        comparison_candidates=scalar(
            "SELECT COUNT(*) FROM fact_dimension_comparison_candidate"
        ),
        combination_title_candidates=scalar(
            """
            SELECT COUNT(*)
            FROM fact_dimension_resolution_ledger
            WHERE product_name LIKE '%세트%'
               OR product_name LIKE '%패키지%'
               OR lower(product_name) LIKE '%package%'
               OR instr(product_name, '+') > 0
               OR instr(product_name, '＋') > 0
            """
        ),
        combination_high_confidence=scalar(
            """
            SELECT COUNT(*)
            FROM fact_dimension_resolution_ledger
            WHERE product_name LIKE '%세트%'
               OR product_name LIKE '%패키지%'
               OR lower(product_name) LIKE '%package%'
               OR (
                    mid_category='소파'
                    AND (
                        instr(product_name, '+') > 0
                        OR instr(product_name, '＋') > 0
                    )
                    AND (
                        product_name LIKE '%스툴%'
                        OR lower(product_name) LIKE '%ottoman%'
                        OR product_name LIKE '%오토만%'
                    )
               )
            """
        ),
        sofa_stool_combinations=scalar(
            """
            SELECT COUNT(*)
            FROM fact_dimension_resolution_ledger
            WHERE mid_category='소파'
              AND (
                  instr(product_name, '+') > 0
                  OR instr(product_name, '＋') > 0
              )
              AND (
                  product_name LIKE '%스툴%'
                  OR lower(product_name) LIKE '%ottoman%'
                  OR product_name LIKE '%오토만%'
              )
            """
        ),
        component_dimensions_fully_confirmed=scalar(
            """
            SELECT COUNT(*)
            FROM vw_product_combination_current
            WHERE component_output_status='ALL_COMPONENT_DIMENSIONS_CONFIRMED'
            """
        ),
        component_rows=int(workbook_meta["component_rows"]),
        bed_single_asset=scalar(
            """
            SELECT COUNT(*)
            FROM vw_bed_asset_policy_current
            WHERE asset_policy='SINGLE_3D_ASSET'
            """
        ),
        workbook_columns=int(workbook_meta["fields"]) + 3,
        workbook_sheets=9,
        pattern_no_check=int(
            workbook_meta["pattern_status_counts"]["완료_패턴체크불필요"]
        ),
        pattern_dimension=int(
            workbook_meta["pattern_status_counts"]["체크필요_규격패턴"]
        ),
        pattern_combination=int(
            workbook_meta["pattern_status_counts"]["체크필요_조합패턴"]
        ),
        pattern_both=int(
            workbook_meta["pattern_status_counts"]["체크필요_규격+조합패턴"]
        ),
        dimension_patterns={
            key: int(value)
            for key, value in workbook_meta["dimension_pattern_counts"].items()
        },
        combination_pattern_groups={
            key: int(value)
            for key, value in workbook_meta[
                "combination_pattern_group_counts"
            ].items()
        },
        remaining=remaining,
    )
    conn.close()
    return stats


def set_fill(shape, color: RGBColor, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency


def set_line(shape, color: RGBColor = LINE, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 16,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.02,
    font: str = FONT,
    linespacing: float | None = None,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if linespacing:
        p.line_spacing = linespacing
    return box


def add_rich_text(
    slide,
    runs: Sequence[tuple[str, float, RGBColor, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for value, size, color, bold in runs:
        run = p.add_run()
        run.text = value
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor = WHITE,
    line: RGBColor = LINE,
    radius: bool = True,
    shadow: bool = False,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    set_line(shape, line, 0.8)
    if shadow:
        shape.shadow.inherit = False
        shape.shadow.blur_radius = Pt(5)
        shape.shadow.distance = Pt(2)
        shape.shadow.transparency = 78
    return shape


def add_pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    *,
    fill: RGBColor,
    color: RGBColor,
    size: float = 10,
):
    shape = add_box(slide, x, y, w, 0.34, fill=fill, line=fill)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.28,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return shape


def add_slide_base(prs: Presentation, section: str, title: str, page: int, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = BG
    add_text(slide, section.upper(), 0.68, 0.34, 4.8, 0.25, size=9, color=LG_RED, bold=True)
    add_text(slide, title, 0.68, 0.68, 11.95, 0.55, size=24, color=INK, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.22, 11.9, 0.38, size=10.5, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.68), Inches(1.58), Inches(11.95), Inches(0.012)
    )
    set_fill(line, LINE)
    line.line.fill.background()
    add_text(
        slide,
        f"{page:02d}",
        12.37,
        7.08,
        0.28,
        0.2,
        size=8,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        f"HomeStyle product data PoC · {REPORT_DATE}",
        0.7,
        7.08,
        4.1,
        0.2,
        size=7.5,
        color=MUTED,
    )
    return slide


def add_bullet_list(
    slide,
    items: Iterable[str],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 13,
    color: RGBColor = INK,
    bullet_color: RGBColor = LG_RED,
    gap: float = 0.42,
):
    for idx, item in enumerate(items):
        yy = y + idx * gap
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(yy + 0.1), Inches(0.08), Inches(0.08)
        )
        set_fill(dot, bullet_color)
        dot.line.fill.background()
        add_text(slide, item, x + 0.18, yy, w - 0.18, gap, size=size, color=color)


def add_kpi_card(
    slide,
    x: float,
    y: float,
    w: float,
    title: str,
    value: str,
    note: str,
    accent: RGBColor,
    light: RGBColor,
):
    add_box(slide, x, y, w, 1.42, fill=WHITE, line=LINE, shadow=True)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(1.42)
    )
    set_fill(accent_bar, accent)
    accent_bar.line.fill.background()
    add_pill(slide, title, x + 0.25, y + 0.18, min(w - 0.48, 1.32), fill=light, color=accent)
    add_text(slide, value, x + 0.25, y + 0.56, w - 0.48, 0.48, size=25, color=INK, bold=True)
    add_text(slide, note, x + 0.25, y + 1.07, w - 0.48, 0.24, size=8.5, color=MUTED)


def add_table(
    slide,
    headers: Sequence[str],
    rows: Sequence[Sequence[str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    col_widths: Sequence[float] | None = None,
    header_fill: RGBColor = INK,
    header_color: RGBColor = WHITE,
    font_size: float = 9.2,
    first_col_bold: bool = False,
    row_fills: Sequence[RGBColor] | None = None,
):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.035)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.font.name = FONT
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = header_color
    for row_idx, row in enumerate(rows, start=1):
        fill_color = row_fills[row_idx - 1] if row_fills else (WHITE if row_idx % 2 else BG)
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.035)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = FONT
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = INK
                paragraph.font.bold = first_col_bold and col_idx == 0
        for col_idx in range(len(headers)):
            cell = table.cell(row_idx, col_idx)
            cell.border_top = None
    return table_shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor = MUTED):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(1.6)
    line.line.end_arrowhead = True
    return line


def add_stage(
    slide,
    x: float,
    y: float,
    w: float,
    number: str,
    title: str,
    body: str,
    accent: RGBColor,
    light: RGBColor,
):
    add_box(slide, x, y, w, 1.35, fill=WHITE, line=LINE)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.17), Inches(0.42), Inches(0.42)
    )
    set_fill(circle, accent)
    circle.line.fill.background()
    add_text(
        slide,
        number,
        x + 0.18,
        y + 0.18,
        0.42,
        0.36,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, title, x + 0.72, y + 0.17, w - 0.9, 0.32, size=12, bold=True)
    add_text(slide, body, x + 0.18, y + 0.68, w - 0.36, 0.5, size=9.3, color=MUTED)
    tag = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 1.28), Inches(w), Inches(0.07)
    )
    set_fill(tag, accent)
    tag.line.fill.background()


def make_cover(prs: Presentation, stats: Stats):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = INK
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(SLIDE_H)
    )
    set_fill(block, LG_RED)
    block.line.fill.background()

    add_pill(slide, "DEVELOPER HANDOFF · POC REPORT", 0.85, 0.65, 2.65, fill=LG_RED_DARK, color=WHITE, size=9)
    add_text(
        slide,
        "홈스타일 상품정보 수집·규격 OCR\n개발 인수인계 보고",
        0.85,
        1.32,
        8.5,
        1.58,
        size=31,
        color=WHITE,
        bold=True,
        linespacing=1.02,
    )
    add_text(
        slide,
        "API · HTML · FAQ/Q&A · 2단계 이미지 OCR 기반\n요청 1·2 필드 산출 및 제품화 준비 상태",
        0.88,
        3.18,
        7.5,
        0.9,
        size=15,
        color=rgb("#D9DEE5"),
    )

    add_box(slide, 9.55, 0.78, 2.75, 5.68, fill=rgb("#2B2D31"), line=rgb("#404348"))
    add_text(slide, "CURRENT SNAPSHOT", 9.88, 1.12, 2.0, 0.25, size=9, color=rgb("#B9C0C8"), bold=True)
    add_text(slide, f"{stats.total:,}", 9.85, 1.56, 2.0, 0.52, size=30, color=WHITE, bold=True)
    add_text(slide, "홈스타일 상품", 9.88, 2.08, 1.8, 0.24, size=10, color=rgb("#B9C0C8"))
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.88), Inches(2.58), Inches(1.98), Inches(0.012)
    )
    set_fill(line, rgb("#4C5057"))
    line.line.fill.background()
    add_text(slide, f"{stats.locked:,}", 9.85, 2.9, 2.0, 0.48, size=27, color=rgb("#8AB4F8"), bold=True)
    add_text(slide, "사용 규격 확정값", 9.88, 3.4, 1.8, 0.24, size=10, color=rgb("#B9C0C8"))
    add_text(slide, f"{stats.pattern_check_needed:,}", 9.85, 4.14, 2.0, 0.48, size=27, color=rgb("#FDD663"), bold=True)
    add_text(slide, "패턴 체크 필요", 9.88, 4.64, 1.8, 0.24, size=10, color=rgb("#B9C0C8"))
    add_text(slide, f"{stats.no_candidate:,}", 9.85, 5.38, 2.0, 0.48, size=27, color=rgb("#F28B82"), bold=True)
    add_text(slide, "최종 규격 무후보", 9.88, 5.88, 1.8, 0.24, size=10, color=rgb("#B9C0C8"))

    add_text(slide, "2026. 07. 27", 0.88, 6.82, 1.5, 0.24, size=9, color=rgb("#AEB5BD"))
    add_text(slide, "LG HomeStyle Data PoC", 2.48, 6.82, 2.5, 0.24, size=9, color=rgb("#AEB5BD"))


def slide_executive(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Executive summary",
        "결론: 9,358개 산출 상태는 모두 분류 완료, 패턴 체크 3,484개를 별도 관리",
        page,
        "‘산출 완료’, ‘패턴 체크’, ‘필수값 PASS’는 서로 다른 축이며 합치거나 대체해서 해석하지 않습니다.",
    )
    card_w = 2.8
    x_positions = [0.7, 3.74, 6.78, 9.82]
    add_kpi_card(slide, x_positions[0], 1.88, card_w, "대상", f"{stats.total:,}개", f"홈스타일 비음영 · {stats.categories}개 카테고리", LG_RED, LG_PINK)
    add_kpi_card(slide, x_positions[1], 1.88, card_w, "산출 상태 분류", "100.0%", f"확정 {stats.locked:,} + RAW {stats.comparison:,} + 무후보 {stats.no_candidate:,}", BLUE, BLUE_LIGHT)
    add_kpi_card(slide, x_positions[2], 1.88, card_w, "패턴 체크 필요", f"{stats.pattern_check_needed:,}개", f"규격 {stats.pattern_dimension:,} · 조합 {stats.pattern_combination:,} · 동시 {stats.pattern_both:,}", AMBER, AMBER_LIGHT)
    add_kpi_card(slide, x_positions[3], 1.88, card_w, "별첨 필수값 PASS", f"{stats.pass_pct:.1f}%", f"{stats.mandatory_pass:,} PASS / {stats.reinforcement:,} 보강", GREEN, GREEN_LIGHT)

    add_box(slide, 0.7, 3.62, 7.5, 2.65, fill=WHITE, line=LINE)
    add_text(slide, "현재 확인된 것", 0.98, 3.91, 2.1, 0.32, size=15, bold=True)
    add_bullet_list(
        slide,
        [
            "API·HTML·Q&A·OCR를 함께 사용해야 고객 요청 필드와 규격을 충분히 채울 수 있음",
            f"요청 1·2를 {stats.workbook_columns}열 Excel로 산출하고 상태·패턴을 상품 ID별로 필터링 가능",
            "LLM 없이 정규식·카테고리·옵션·문맥 규칙으로 재현 가능한 결과를 생성함",
            "위험한 값은 RAW 비교정보로 보존하고 D00~D05·N/S/Q/M/A/O 패턴으로 검토 대상을 축소함",
        ],
        1.0,
        4.38,
        6.9,
        1.7,
        size=11.2,
        bullet_color=GREEN,
        gap=0.39,
    )
    add_box(slide, 8.47, 3.62, 4.16, 2.65, fill=AMBER_LIGHT, line=rgb("#F2C36B"))
    add_text(slide, "제품화 전 핵심 조치", 8.76, 3.91, 2.7, 0.32, size=15, bold=True, color=AMBER)
    add_bullet_list(
        slide,
        [
            "빈 DB bootstrap 순환 의존성 해소",
            "단일 실행기·설정·requirements 추가",
            "패턴 승인 결과의 DB 상태 전환·회귀 테스트",
            "scope 삭제 동기화·run ID·운영 로그",
        ],
        8.78,
        4.39,
        3.5,
        1.62,
        size=11,
        bullet_color=AMBER,
        gap=0.39,
    )


def slide_requests(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Scope",
        "고객 요청 1·2는 모두 포함했고, 홈스타일 조건에 맞게 적용 여부를 분리",
        page,
        "가전 전용 설치 타입과 개인 CRM은 현재 홈스타일 공개 데이터 범위에서 ‘해당없음/미확보’로 명시합니다.",
    )
    add_box(slide, 0.7, 1.86, 5.95, 4.83, fill=WHITE, line=BLUE)
    add_pill(slide, "요청 1 · 3D Asset 기본 정보", 1.0, 2.12, 2.5, fill=BLUE_LIGHT, color=BLUE, size=10.5)
    req1 = [
        ("필수", "PDP 이미지 · 롤링 이미지 URL"),
        ("필수", "W / D / H 숫자 규격(mm)"),
        ("필수", "중카테고리 · 소카테고리"),
        ("필수", "추천 공간 · 브랜드 · 색상"),
        ("조건부", "가전 설치 타입 · 세트 구성 ID"),
        ("옵션", "배치 위치 · 벽면 추천 높이"),
        ("확장", "옵션 스타일 · 옵션 01…N"),
    ]
    for idx, (tag, item) in enumerate(req1):
        y = 2.7 + idx * 0.49
        fill = BLUE_LIGHT if tag == "필수" else GRAY_LIGHT
        color = BLUE if tag == "필수" else MUTED
        add_pill(slide, tag, 1.0, y, 0.64, fill=fill, color=color, size=8.5)
        add_text(slide, item, 1.79, y + 0.02, 4.38, 0.28, size=11.2, color=INK)
    add_text(
        slide,
        "별첨 0715 PASS 기준: 브랜드·분류·ID·이미지·색상·완성 W/D/H\n세트 실제 ID는 합의에 따라 필수 판정에서 제외",
        1.0,
        6.18,
        5.15,
        0.42,
        size=8.7,
        color=MUTED,
    )

    add_box(slide, 6.9, 1.86, 5.73, 4.83, fill=WHITE, line=GREEN)
    add_pill(slide, "요청 2 · 추천 분석 활용", 7.2, 2.12, 2.25, fill=GREEN_LIGHT, color=GREEN, size=10.5)
    req2 = [
        ("01", "설명서 13개 영역 자동 태깅"),
        ("02", "디자인 스타일 규칙 추론"),
        ("03", "공간 콘텐츠 자동 태깅"),
        ("04", "공간스타일 ↔ 제품 관계"),
        ("05", "공간 내 제품 ↔ 제품 관계"),
        ("06", "개인 구매·선호·CRM"),
        ("07", "색상 동의어 기반 검색 문자열"),
    ]
    for idx, (num, item) in enumerate(req2):
        y = 2.7 + idx * 0.49
        add_pill(slide, num, 7.2, y, 0.52, fill=GREEN_LIGHT, color=GREEN, size=8.5)
        add_text(slide, item, 7.9, y + 0.02, 4.1, 0.28, size=11.2)
    add_text(
        slide,
        "현재는 규칙 기반 태그·관계 문자열 산출 단계\n임베딩 검색·CRM 결합·관계 그래프 서비스는 후속 개발 범위",
        7.2,
        6.18,
        4.85,
        0.42,
        size=8.7,
        color=MUTED,
    )


def slide_input_output(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Operating model",
        "입력은 카테고리 scope 또는 상품 ID, 출력은 증거를 가진 DB와 요청 필드 Excel",
        page,
        "현재 scope는 제품군 Excel에서 만든 cache를 사용하며, URL은 productId로 동일하게 연결됩니다.",
    )
    stages = [
        ("01", "대상 입력", "비음영 카테고리 scope\n또는 productId 목록", LG_RED, LG_PINK),
        ("02", "원천 수집", "상품·공간·세트·Q&A API\nPDP HTML·상세 이미지", BLUE, BLUE_LIGHT),
        ("03", "해석·OCR", "기본 파싱 → 1차 스캔\n→ layout/영역 재OCR", TEAL, TEAL_LIGHT),
        ("04", "상태 판정", "산출 상태 + 패턴 상태\nD / N·S·Q·M·A·O", AMBER, AMBER_LIGHT),
        ("05", "산출", "SQLite 원장·이력\n9시트 Excel·필터", GREEN, GREEN_LIGHT),
    ]
    x_values = [0.7, 3.19, 5.68, 8.17, 10.66]
    for idx, (num, title, body, accent, light) in enumerate(stages):
        add_stage(slide, x_values[idx], 2.18, 1.98, num, title, body, accent, light)
        if idx < len(stages) - 1:
            add_arrow(slide, x_values[idx] + 2.02, 2.85, x_values[idx + 1] - 0.04, 2.85)

    add_box(slide, 0.7, 4.16, 5.95, 1.83, fill=WHITE, line=LINE)
    add_text(slide, "입력 예", 0.98, 4.44, 1.2, 0.28, size=13, bold=True)
    add_text(slide, "productId", 0.98, 4.96, 1.0, 0.24, size=9, color=MUTED, bold=True)
    add_text(slide, "G25070005743", 1.92, 4.91, 2.3, 0.34, size=15, color=LG_RED, bold=True)
    add_text(slide, "URL", 0.98, 5.43, 0.8, 0.22, size=9, color=MUTED, bold=True)
    add_text(slide, "homestyle.lge.co.kr/item?productId=…", 1.92, 5.37, 4.25, 0.32, size=10.5, color=BLUE)

    add_box(slide, 6.9, 4.16, 5.73, 1.83, fill=WHITE, line=LINE)
    add_text(slide, "출력 예", 7.18, 4.44, 1.2, 0.28, size=13, bold=True)
    add_text(slide, "대표 W/D/H", 7.18, 4.96, 1.15, 0.24, size=9, color=MUTED, bold=True)
    add_text(slide, "2910 / 1020 / 910 mm", 8.37, 4.91, 2.9, 0.34, size=15, color=GREEN, bold=True)
    add_text(slide, "근거", 7.18, 5.43, 0.8, 0.22, size=9, color=MUTED, bold=True)
    add_text(slide, "상품 API 고시정보 · SOURCE_CONFIRMED", 8.37, 5.38, 3.75, 0.3, size=10.5, color=INK)

    add_pill(slide, f"DB {stats.total:,}행 원장", 3.36, 6.34, 1.72, fill=GRAY_LIGHT, color=INK, size=9)
    add_pill(
        slide,
        f"Excel {stats.workbook_sheets}개 시트 · {stats.workbook_columns}열",
        5.28,
        6.34,
        2.25,
        fill=GRAY_LIGHT,
        color=INK,
        size=9,
    )
    add_pill(slide, "필드별 값·상태·비교근거", 7.58, 6.34, 2.15, fill=GRAY_LIGHT, color=INK, size=9)


def slide_sources(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Source depth",
        "API 단독으로 끝내지 않고 HTML·FAQ/Q&A·이미지를 동일 상품 ID로 결합",
        page,
        "현재 9,358개 상품에서 상품·공간·세트·Q&A API와 PDP HTML 응답은 모두 HTTP 200으로 저장되어 있습니다.",
    )
    headers = ["원천", "현재 상태", "주요 산출", "한계 / 보완"]
    rows = [
        ["상품 상세 API", f"{stats.sources:,} / {stats.total:,}", "상품명·브랜드·이미지·옵션·고시", "판매처 표기가 ‘상세 참조’인 경우"],
        ["공간 추천 API", f"{stats.sources:,} / {stats.total:,}", "공간·스타일·연관 제품", "없는 값은 카테고리 fallback"],
        ["세트 API", f"{stats.sources:,} / {stats.total:,}", "구성 productId", "명칭만 세트인 경우 실제 ID 미확보"],
        ["Q&A API", f"{stats.sources:,} / {stats.total:,}", "재질·규격·설치 보강", "상세 본문 추가 조회는 상품당 최대 3건"],
        ["PDP HTML", f"{stats.sources:,} / {stats.total:,}", "JSON-LD FAQ·visible text·이미지 URL", "동적 상세는 이미지 OCR 필요"],
        ["상세 이미지 OCR", "필요 상품 선별", "규격·재질·주의·INFO/Dimension", "해상도·배치·문맥 오탐 관리"],
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.93,
        11.93,
        3.72,
        col_widths=[1.85, 1.45, 3.2, 5.43],
        header_fill=INK,
        font_size=9.4,
        first_col_bold=True,
    )
    add_box(slide, 0.7, 5.94, 11.93, 0.78, fill=BLUE_LIGHT, line=rgb("#C7D9FA"))
    add_rich_text(
        slide,
        [
            ("핵심 설계 원칙  ", 11, BLUE_DARK, True),
            ("상품 API 명시값 → HTML/FAQ·Q&A → 검증된 OCR → 카테고리 추론 → 미확보/해당없음", 11, INK, False),
        ],
        1.0,
        6.18,
        11.2,
        0.28,
        valign=MSO_ANCHOR.MIDDLE,
    )


def slide_ocr(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "OCR pipeline",
        "전체 이미지 스캔 후 후보 영역만 고해상도 OCR하는 2단계 구조",
        page,
        "비용을 줄이면서도 INFO·SIZE·DIMENSION·제품 사이즈 영역의 작은 글자를 다시 읽도록 설계했습니다.",
    )
    data = [
        ("1차 전체 스캔", f"{stats.pass1_images:,}", "상품-이미지", f"{stats.pass1_products:,}개 상품", LG_RED, LG_PINK),
        ("2차 Layout OCR", f"{stats.pass2_observations:,}", "좌표 관측", f"{stats.pass2_products:,}개 상품", BLUE, BLUE_LIGHT),
        ("제목 영역 재OCR", "2,659", "최종 캠페인 영역", "2,659건 OCR 성공", TEAL, TEAL_LIGHT),
        ("비교 후보 축적", f"{stats.comparison_candidates:,}", "후보 행", f"{stats.comparison:,}개 상품 제공", GREEN, GREEN_LIGHT),
    ]
    for idx, (title, value, unit, note, color, light) in enumerate(data):
        x = 0.7 + idx * 3.04
        add_box(slide, x, 1.92, 2.8, 1.72, fill=WHITE, line=color)
        add_pill(slide, title, x + 0.22, 2.14, 1.54, fill=light, color=color, size=8.8)
        add_text(slide, value, x + 0.22, 2.6, 1.45, 0.42, size=23, bold=True, color=INK)
        add_text(slide, unit, x + 1.52, 2.73, 1.0, 0.22, size=8.5, color=MUTED)
        add_text(slide, note, x + 0.22, 3.18, 2.24, 0.22, size=8.3, color=MUTED)

    add_box(slide, 0.7, 4.0, 11.93, 2.33, fill=WHITE, line=LINE)
    step_x = [1.0, 3.45, 5.9, 8.35, 10.8]
    step_titles = ["URL 목록화", "저비용 OCR", "후보 점수화", "영역 crop", "좌표·문맥 파싱"]
    step_notes = [
        "상세 이미지 전체",
        "SIZE·단위·숫자",
        "배송 배너 감점",
        "제목 위/아래 확장",
        "축·단위·역할 결합",
    ]
    step_colors = [LG_RED, LG_RED, BLUE, TEAL, GREEN]
    for idx in range(5):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(step_x[idx]), Inches(4.49), Inches(0.58), Inches(0.58)
        )
        set_fill(circle, step_colors[idx])
        circle.line.fill.background()
        add_text(
            slide,
            str(idx + 1),
            step_x[idx],
            4.52,
            0.58,
            0.38,
            size=12,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, step_titles[idx], step_x[idx] - 0.25, 5.2, 1.15, 0.26, size=10.2, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, step_notes[idx], step_x[idx] - 0.42, 5.57, 1.48, 0.38, size=8.3, color=MUTED, align=PP_ALIGN.CENTER)
        if idx < 4:
            add_arrow(slide, step_x[idx] + 0.72, 4.78, step_x[idx + 1] - 0.12, 4.78, LINE)
    add_text(
        slide,
        f"참고: DB에는 targeted crop 누적 {stats.targeted_crop_total:,}건이 보존되어 있으며, 2,659건은 최종 잔여 대상 재OCR 캠페인 수치입니다.",
        0.72,
        6.55,
        11.2,
        0.22,
        size=7.8,
        color=MUTED,
    )


def slide_rules(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Normalization",
        "확정 가능한 규칙과 자동 확정 금지 규칙을 분리해 오탐을 제어",
        page,
        "값을 많이 채우는 것보다 제품 대표 규격인지 구분하는 것이 핵심입니다.",
    )
    add_box(slide, 0.7, 1.9, 5.8, 4.82, fill=GREEN_LIGHT, line=rgb("#A8DAB5"))
    add_pill(slide, "AUTO-LOCK", 1.0, 2.16, 1.1, fill=GREEN, color=WHITE, size=9)
    add_text(slide, "자동 확정 가능", 1.0, 2.66, 2.4, 0.35, size=18, color=GREEN, bold=True)
    add_bullet_list(
        slide,
        [
            "명시적 W/D/H, 가로·깊이·높이",
            "가구 문맥의 순서형 3개 값",
            "L/D/H · W/L/H · L/W/H → W/D/H",
            "원형 지름+높이 → W=D=지름",
            "2D 액자·포스터 → D 비적용",
            "상품명·모델명·옵션 코드 일치",
            "M/L, SS/Q/K 등 복수 옵션 분리",
        ],
        1.03,
        3.18,
        4.92,
        2.75,
        size=11.2,
        bullet_color=GREEN,
        gap=0.38,
    )
    add_text(slide, "단위: mm 그대로 · cm ×10 · inch ×25.4", 1.03, 6.27, 4.6, 0.24, size=9, color=MUTED)

    add_box(slide, 6.83, 1.9, 5.8, 4.82, fill=RED_LIGHT, line=rgb("#F0B4B2"))
    add_pill(slide, "BLOCK / COMPARE", 7.13, 2.16, 1.48, fill=RED, color=WHITE, size=9)
    add_text(slide, "자동 확정 금지", 7.13, 2.66, 2.4, 0.35, size=18, color=RED, bold=True)
    add_bullet_list(
        slide,
        [
            "배송·반입·엘리베이터·사다리차 수치",
            "가드·발통·선반 등 구성품 수치",
            "다른 라인업 모델의 규격",
            "허용오차 ±3cm",
            "축 범위가 모호한 일부 값",
            "pass2 원시 숫자 순서·추정 단위",
            "충돌하는 높이·옵션은 비교 후보 제공",
        ],
        7.16,
        3.18,
        4.92,
        2.75,
        size=11.2,
        bullet_color=RED,
        gap=0.38,
    )
    add_text(slide, "원칙: 오탐 후보는 폐기하지 않고 원문·대상·축과 함께 분리 저장", 7.16, 6.27, 4.85, 0.24, size=9, color=MUTED)


def slide_dimension_result(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Dimension output status",
        f"산출 상태는 전 상품 처리 결과: 확정값 {stats.locked:,} · RAW 분리 {stats.comparison:,} · 무후보 분류 {stats.no_candidate:,}",
        page,
        "모든 상품이 세 상태 중 하나를 가지므로 산출 처리율은 100%입니다. 3D 입력값 확정률과는 다른 개념입니다.",
    )
    x, y, w, h = 0.7, 2.12, 11.93, 0.72
    segments = [
        ("확정값·API/HTML", stats.source_confirmed, LG_RED),
        ("확정값·OCR/규칙", stats.rule_resolved, BLUE),
        ("완료_RAW분리", stats.comparison, AMBER),
        ("완료_무후보분류", stats.no_candidate, RED),
    ]
    current_x = x
    for label, value, color in segments:
        seg_w = w * value / stats.total
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(current_x),
            Inches(y),
            Inches(seg_w),
            Inches(h),
        )
        set_fill(shape, color)
        shape.line.fill.background()
        current_x += seg_w
    add_text(slide, "0", x, 2.92, 0.5, 0.2, size=8, color=MUTED)
    add_text(slide, f"{stats.total:,}", x + w - 0.7, 2.92, 0.7, 0.2, size=8, color=MUTED, align=PP_ALIGN.RIGHT)

    for idx, (label, value, color) in enumerate(segments):
        xx = 0.7 + idx * 3.04
        add_box(slide, xx, 3.42, 2.8, 1.26, fill=WHITE, line=LINE)
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(xx + 0.22), Inches(3.7), Inches(0.18), Inches(0.18)
        )
        set_fill(dot, color)
        dot.line.fill.background()
        add_text(slide, label, xx + 0.5, 3.62, 1.92, 0.28, size=10, bold=True)
        add_text(slide, f"{value:,}", xx + 0.22, 4.03, 1.1, 0.38, size=21, bold=True, color=color)
        add_text(slide, f"{value / stats.total * 100:.1f}%", xx + 1.48, 4.12, 0.75, 0.24, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    add_box(slide, 0.7, 5.14, 7.58, 1.32, fill=BLUE_LIGHT, line=rgb("#C7D9FA"))
    add_text(slide, "산출 처리 완료", 1.0, 5.47, 1.55, 0.26, size=11, color=BLUE_DARK, bold=True)
    add_text(slide, f"{stats.total:,}", 2.7, 5.36, 1.6, 0.44, size=25, color=BLUE, bold=True)
    add_text(slide, "100.0%", 4.25, 5.43, 1.0, 0.32, size=16, color=BLUE_DARK, bold=True)
    add_text(slide, "= 확정값 + RAW 비교정보 + 무후보 원인분류", 5.25, 5.48, 2.75, 0.24, size=9.0, color=MUTED)
    add_box(slide, 8.55, 5.14, 4.08, 1.32, fill=AMBER_LIGHT, line=rgb("#F2C36B"))
    add_text(slide, "별도 패턴 체크 축", 8.86, 5.4, 1.7, 0.22, size=9, color=AMBER, bold=True)
    add_text(slide, f"{stats.pattern_check_needed:,}개", 8.86, 5.77, 1.35, 0.34, size=19, color=AMBER, bold=True)
    add_text(slide, "규격·조합 규칙 승인 대상", 10.2, 5.83, 1.9, 0.24, size=8.2, color=INK)


def slide_mandatory(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Mandatory PASS",
        f"별첨 필수값 6개 그룹을 모두 충족해야 PASS — 현재 {stats.mandatory_pass:,}개, {stats.pass_pct:.1f}%",
        page,
        "세트 구성 실제 ID는 합의에 따라 필수 판정에서 제외했으며, 가전 조건부 필드는 홈스타일에 적용하지 않습니다.",
    )
    add_box(slide, 0.7, 1.92, 5.55, 4.7, fill=WHITE, line=LINE)
    add_text(slide, "상품 판정", 1.0, 2.23, 1.5, 0.28, size=14, bold=True)
    total_w = 4.95
    pass_w = total_w * stats.mandatory_pass / stats.total
    fail_w = total_w - pass_w
    pbar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.87), Inches(pass_w), Inches(0.65)
    )
    set_fill(pbar, GREEN)
    pbar.line.fill.background()
    fbar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.0 + pass_w), Inches(2.87), Inches(fail_w), Inches(0.65)
    )
    set_fill(fbar, RED)
    fbar.line.fill.background()
    add_text(slide, "PASS", 1.0, 3.72, 0.8, 0.24, size=9, color=GREEN, bold=True)
    add_text(slide, f"{stats.mandatory_pass:,}", 1.0, 4.03, 1.65, 0.44, size=26, color=GREEN, bold=True)
    add_text(slide, f"{stats.pass_pct:.1f}%", 2.54, 4.14, 0.8, 0.24, size=11, color=MUTED)
    add_text(slide, "보강대상", 3.65, 3.72, 1.0, 0.24, size=9, color=RED, bold=True)
    add_text(slide, f"{stats.reinforcement:,}", 3.65, 4.03, 1.65, 0.44, size=26, color=RED, bold=True)
    add_text(slide, f"{stats.reinforcement / stats.total * 100:.1f}%", 5.05, 4.14, 0.72, 0.24, size=11, color=MUTED)
    add_text(
        slide,
        "PASS = 브랜드 + 중·소분류 + 상품 ID + 이미지\n       + 색상 + 완성 W/D/H",
        1.0,
        5.02,
        4.7,
        0.72,
        size=10.5,
        color=INK,
    )

    add_box(slide, 6.55, 1.92, 6.08, 4.7, fill=WHITE, line=LINE)
    add_text(slide, "주요 보강 필드", 6.86, 2.23, 1.8, 0.28, size=14, bold=True)
    missing_rows = [
        ("완성 W/D/H", stats.size_missing, BLUE, "비교 후보라도 축 하나가 없으면 미충족"),
        ("제품 색상", stats.color_missing, LG_RED, "고시·색상 옵션이 없으면 미확보"),
    ]
    max_val = max(v for _, v, _, _ in missing_rows)
    for idx, (label, value, color, note) in enumerate(missing_rows):
        yy = 2.98 + idx * 1.35
        add_text(slide, label, 6.86, yy, 1.25, 0.24, size=10.5, bold=True)
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(8.15), Inches(yy + 0.02), Inches(3.4), Inches(0.32)
        )
        set_fill(bar_bg, GRAY_LIGHT)
        bar_bg.line.fill.background()
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(8.15),
            Inches(yy + 0.02),
            Inches(3.4 * value / max_val),
            Inches(0.32),
        )
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(slide, f"{value:,}", 11.67, yy - 0.01, 0.62, 0.28, size=11, color=color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, note, 8.15, yy + 0.47, 3.85, 0.28, size=8.4, color=MUTED)
    add_box(slide, 6.86, 5.72, 5.15, 0.56, fill=AMBER_LIGHT, line=rgb("#F2C36B"))
    add_text(
        slide,
        f"주의: 최종 무후보 {stats.no_candidate:,}개와 필수 사이즈 미충족 {stats.size_missing:,}개는 다른 통계",
        7.08,
        5.89,
        4.75,
        0.21,
        size=9.2,
        color=AMBER,
        bold=True,
    )


def slide_remaining(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Remaining backlog",
        f"최종 무후보 {stats.no_candidate:,}개의 65%는 전체 이미지에서 크기 신호 자체가 없음",
        page,
        "남은 문제는 하나가 아니라 ‘원천 부재·영역 탐지·OCR 인식·수집 오류’로 분리해 처리해야 합니다.",
    )
    max_count = max(count for _, _, count in stats.remaining)
    for idx, (code, name, count) in enumerate(stats.remaining):
        yy = 1.95 + idx * 0.69
        add_text(slide, name, 0.72, yy, 2.95, 0.25, size=9.5, bold=True)
        add_text(slide, code, 0.72, yy + 0.29, 2.95, 0.18, size=6.8, color=MUTED)
        bg_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3.83), Inches(yy + 0.06), Inches(6.75), Inches(0.32)
        )
        set_fill(bg_bar, GRAY_LIGHT)
        bg_bar.line.fill.background()
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(3.83),
            Inches(yy + 0.06),
            Inches(6.75 * count / max_count),
            Inches(0.32),
        )
        color = LG_RED if idx == 0 else (AMBER if idx < 3 else MUTED)
        set_fill(bar, color)
        bar.line.fill.background()
        add_text(slide, f"{count:,}", 10.78, yy + 0.01, 0.55, 0.28, size=10.5, color=color, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, f"{count / stats.no_candidate * 100:.1f}%", 11.45, yy + 0.01, 0.68, 0.28, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)

    add_box(slide, 0.7, 6.1, 11.93, 0.58, fill=WHITE, line=LINE)
    add_rich_text(
        slide,
        [
            ("다음 처리 순서  ", 10, LG_RED, True),
            ("① 제조사/API 추가 원천  →  ② 제목 없는 도면 탐지  →  ③ OCR 엔진·전처리 비교  →  ④ 재다운로드·예외 규칙", 10, INK, False),
        ],
        0.98,
        6.28,
        11.15,
        0.22,
    )


def slide_examples(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Examples",
        "원천과 후보 구조에 따라 산출 상태·패턴 상태가 함께 결정",
        page,
        "자동화가 단순 OCR 숫자 추출이 아니라 상품명·옵션·문맥·제외 규칙을 함께 판단하는 이유입니다.",
    )
    examples = [
        (
            "G25070005743",
            "API 직접 확정",
            "[4인] W2910×D1020×H910mm\n[스툴] W740×D660×H410mm",
            "대표 2910 / 1020 / 910\n완료_확정값",
            GREEN,
            GREEN_LIGHT,
        ),
        (
            "G25090019410",
            "OCR 보정·규칙 확정",
            "INFO … 가로기600 × 세로2140\n× 헤드높이800mm",
            "가로기600 → 가로2600\n완료_확정값",
            BLUE,
            BLUE_LIGHT,
        ),
        (
            "G25110025989",
            "RAW 비교정보 제공",
            "W590 H810\nD570 H450",
            "590/570/450, 590/570/810\n완료_RAW분리 · D03",
            AMBER,
            AMBER_LIGHT,
        ),
        (
            "G25070000483",
            "무후보 원인분류",
            "옵션 Ø100, Ø120\n높이·추가 규격 신호 없음",
            "W/D/H 공란\n완료_무후보분류",
            RED,
            RED_LIGHT,
        ),
    ]
    positions = [(0.7, 1.9), (6.78, 1.9), (0.7, 4.33), (6.78, 4.33)]
    for (pid, title, source, output, color, light), (x, y) in zip(examples, positions):
        add_box(slide, x, y, 5.85, 2.05, fill=WHITE, line=color)
        add_pill(slide, pid, x + 0.25, y + 0.21, 1.3, fill=light, color=color, size=8.5)
        add_text(slide, title, x + 1.73, y + 0.23, 2.8, 0.28, size=12, bold=True)
        add_text(slide, "입력", x + 0.25, y + 0.77, 0.48, 0.2, size=8, color=MUTED, bold=True)
        add_text(slide, source, x + 0.75, y + 0.7, 2.2, 0.78, size=9.2, color=INK)
        add_arrow(slide, x + 3.02, y + 1.15, x + 3.46, y + 1.15, color)
        add_text(slide, "산출", x + 3.58, y + 0.77, 0.48, 0.2, size=8, color=MUTED, bold=True)
        add_text(slide, output, x + 4.08, y + 0.7, 1.48, 0.8, size=9.2, color=color, bold=True)
        add_text(slide, "제품 상세 URL은 productId로 동일하게 재현 가능", x + 0.25, y + 1.71, 4.9, 0.18, size=7.3, color=MUTED)


def slide_workbook(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Deliverables",
        "Excel은 보고·검수용, SQLite는 원천·후보·판정·이력을 보존하는 기준 저장소",
        page,
        "상품별 요구필드는 한 행으로 유지하고, 규격 후보·옵션·조합 구성품은 별도 상세 시트로 추적합니다.",
    )
    sheets = [
        ("00", "요약", "산출·패턴·PASS·잔여 원인"),
        ("01", "상품별 요구필드", f"요청 1·2+상태 · 9,358행 · {stats.workbook_columns}열"),
        ("02", "규격 상세", "API/HTML/OCR 기본 규격 레코드"),
        ("03", "옵션 상세", "옵션 스타일·값·ID·이미지"),
        ("04", "필수값 판정", "상품별 6개 필수 그룹"),
        ("05", "카테고리 통계", "scope·상품·필수값 통계"),
        ("06", "수집 오류", "원천별 HTTP·수집 상태"),
        ("07", "조합상품 구성품", f"product_id+component_seq · {stats.component_rows:,}행"),
        ("08", "패턴 상태", "산출/패턴/D/N·S·Q·M·A·O 필터"),
    ]
    for idx, (num, name, note) in enumerate(sheets):
        yy = 1.91 + idx * 0.5
        fill = BLUE_LIGHT if idx == 1 else AMBER_LIGHT if idx == 8 else WHITE
        line_color = BLUE if idx == 1 else AMBER if idx == 8 else LINE
        add_box(slide, 0.7, yy, 5.95, 0.39, fill=fill, line=line_color)
        add_pill(slide, num, 0.83, yy + 0.025, 0.46, fill=INK, color=WHITE, size=7.5)
        add_text(slide, name, 1.45, yy + 0.055, 1.55, 0.2, size=9.4, bold=True)
        add_text(slide, note, 2.95, yy + 0.055, 3.35, 0.2, size=8.2, color=MUTED)

    add_box(slide, 6.9, 1.91, 5.73, 4.7, fill=WHITE, line=LINE)
    add_text(slide, "DB 핵심 계층", 7.2, 2.23, 1.65, 0.28, size=14, bold=True)
    db_layers = [
        ("RAW", "products / sources", "API·HTML·Q&A·OCR 원문 BLOB", LG_RED),
        (
            "STAGE",
            "stg_dimension_* + stg_product_*",
            "규격·조합상품·구성품 snapshot",
            BLUE,
        ),
        ("FACT", "fact_dimension_resolution_ledger", "상품별 최종 규격 상태", GREEN),
        ("FACT", "fact_dimension_comparison_candidate", f"비교 후보 {stats.comparison_candidates:,}행", AMBER),
        ("HIST", "hist_dimension_resolution_event", "상태 변경 이력", PURPLE),
    ]
    for idx, (tag, table_name, note, color) in enumerate(db_layers):
        yy = 2.78 + idx * 0.68
        add_pill(slide, tag, 7.2, yy, 0.58, fill=GRAY_LIGHT, color=color, size=7.7)
        add_text(slide, table_name, 7.98, yy + 0.01, 2.8, 0.22, size=9.7, bold=True, color=INK)
        add_text(slide, note, 7.98, yy + 0.29, 3.95, 0.2, size=7.9, color=MUTED)
    add_text(
        slide,
        "운영 보완: Excel 상태값을 DB 최종 field fact로 평면 저장하고\n필드별 source/evidence/confidence와 승인 패턴 버전을 함께 제공",
        7.2,
        6.02,
        4.75,
        0.45,
        size=8.7,
        color=LG_RED,
        bold=True,
    )


def slide_excel_status_model(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Excel status model",
        "Excel은 ‘산출 상태’와 ‘패턴 상태’를 독립 컬럼으로 저장",
        page,
        "산출 상태는 현재 결과의 형태, 패턴 상태는 사람이 규칙을 추가 승인해야 하는지를 뜻합니다.",
    )

    add_box(slide, 0.7, 1.74, 5.82, 2.16, fill=BLUE_LIGHT, line=BLUE)
    add_pill(slide, "축 1 · 산출 상태", 0.98, 1.98, 1.42, fill=BLUE, color=WHITE, size=9)
    add_text(
        slide,
        "현재 데이터를 어떤 형태로 제공했는가?",
        2.58,
        2.03,
        3.35,
        0.23,
        size=10.5,
        color=BLUE_DARK,
        bold=True,
    )
    output_rows = [
        ("완료_확정값", stats.locked, "대표 W/D/H 사용 가능"),
        ("완료_RAW분리", stats.comparison, "후보·원문을 비교 컬럼에 보존"),
        ("완료_무후보분류", stats.no_candidate, "후보 없음과 원인까지 분류"),
    ]
    for idx, (label, count, note) in enumerate(output_rows):
        yy = 2.52 + idx * 0.42
        add_text(slide, label, 1.0, yy, 1.8, 0.22, size=9.1, bold=True)
        add_text(slide, f"{count:,}", 2.86, yy - 0.02, 0.75, 0.25, size=12, color=BLUE_DARK, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, note, 3.78, yy, 2.35, 0.22, size=8.4, color=MUTED)

    add_box(slide, 6.78, 1.74, 5.85, 2.16, fill=AMBER_LIGHT, line=AMBER)
    add_pill(slide, "축 2 · 패턴 상태", 7.06, 1.98, 1.42, fill=AMBER, color=WHITE, size=9)
    add_text(
        slide,
        "추가 패턴 승인이 필요한가?",
        8.67,
        2.03,
        3.25,
        0.23,
        size=10.5,
        color=AMBER,
        bold=True,
    )
    pattern_rows = [
        ("완료_패턴체크불필요", stats.pattern_no_check, "추가 확인 없음"),
        ("체크필요_규격패턴", stats.pattern_dimension, "D01~D05"),
        ("체크필요_조합패턴", stats.pattern_combination, "N/S/Q/M/A/O"),
        ("체크필요_규격+조합패턴", stats.pattern_both, "D00 → 조합 먼저"),
    ]
    for idx, (label, count, note) in enumerate(pattern_rows):
        yy = 2.44 + idx * 0.35
        add_text(slide, label, 7.08, yy, 2.17, 0.2, size=8.25, bold=True)
        add_text(slide, f"{count:,}", 9.26, yy - 0.01, 0.65, 0.22, size=10.5, color=AMBER, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, note, 10.08, yy, 2.0, 0.2, size=8.1, color=MUTED)

    headers = ["산출 상태", "체크불필요", "조합", "규격", "규격+조합", "합계"]
    rows = [
        ["완료_확정값", f"{stats.locked - stats.pattern_combination:,}", f"{stats.pattern_combination:,}", "0", "0", f"{stats.locked:,}"],
        ["완료_RAW분리", "0", "0", f"{stats.pattern_dimension:,}", f"{stats.pattern_both:,}", f"{stats.comparison:,}"],
        ["완료_무후보분류", f"{stats.no_candidate:,}", "0", "0", "0", f"{stats.no_candidate:,}"],
        ["합계", f"{stats.pattern_no_check:,}", f"{stats.pattern_combination:,}", f"{stats.pattern_dimension:,}", f"{stats.pattern_both:,}", f"{stats.total:,}"],
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        4.16,
        11.93,
        1.7,
        col_widths=[2.45, 2.1, 1.55, 1.75, 2.05, 2.03],
        header_fill=INK,
        font_size=8.8,
        first_col_bold=True,
        row_fills=[BLUE_LIGHT, AMBER_LIGHT, RED_LIGHT, GRAY_LIGHT],
    )
    add_box(slide, 0.7, 6.08, 11.93, 0.63, fill=WHITE, line=LINE)
    add_text(slide, "필터 순서", 0.98, 6.27, 0.8, 0.2, size=9, color=LG_RED, bold=True)
    add_text(
        slide,
        "08_패턴_상태 → 패턴 상태=체크필요_* → 규격 패턴 코드(D00~D05) 또는 조합 패턴군(N/S/Q/M/A/O) → 상품 URL 확인",
        1.9,
        6.24,
        10.1,
        0.24,
        size=8.7,
        color=INK,
        bold=True,
    )


def slide_3d_row_policy(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Request 1 · 3D row policy",
        "한 product_id가 몇 개의 3D Asset 행으로 나뉘는지 먼저 결정",
        page,
        "규격값을 선택하기 전에 상품 구조를 판단해야 소파 크기를 스툴 크기와 합치는 오탐을 막을 수 있습니다.",
    )

    add_box(slide, 0.7, 1.82, 2.05, 1.06, fill=INK, line=INK)
    add_text(slide, "INPUT", 0.95, 2.03, 0.62, 0.2, size=8, color=rgb("#B9C0C8"), bold=True)
    add_text(slide, "product_id 1건", 0.95, 2.36, 1.45, 0.3, size=15, color=WHITE, bold=True)
    add_arrow(slide, 2.78, 2.35, 3.32, 2.35, LG_RED)
    add_box(slide, 3.38, 1.82, 3.05, 1.06, fill=LG_PINK, line=LG_RED)
    add_text(slide, "01 · 상품 구조 판정", 3.66, 2.03, 2.2, 0.23, size=10.5, color=LG_RED, bold=True)
    add_text(slide, "상품명 + API 옵션 + 규격 문맥", 3.66, 2.4, 2.35, 0.22, size=8.8, color=INK)
    add_arrow(slide, 6.47, 2.35, 7.02, 2.35, LG_RED)
    add_box(slide, 7.08, 1.82, 5.55, 1.06, fill=WHITE, line=LINE)
    add_text(slide, "02 · N / S / Q / M / A / O 패턴으로 출력 단위 결정", 7.38, 2.04, 4.7, 0.24, size=11, bold=True)
    add_text(slide, "그 다음 각 행에 대표·옵션·구성품 W/D/H를 연결", 7.38, 2.41, 4.35, 0.22, size=8.8, color=MUTED)

    lanes = [
        ("N · S", "단일 Asset 1행", "‘+’ 과탐 또는 조립 후 한 몸\nproduct_id + seq=1", BLUE, BLUE_LIGHT),
        ("Q", "동일 Asset + 수량", "형상은 하나, quantity로 관리\n예: 동일 의자 2P", TEAL, TEAL_LIGHT),
        ("M", "구성품별 복수 행", "서로 독립 배치되는 제품\nproduct_id + component_seq", LG_RED, LG_PINK),
        ("O", "옵션별 복수 행", "외형·규격이 바뀌는 옵션만\nproduct_id + option_seq", AMBER, AMBER_LIGHT),
        ("A", "부속품 정책 판단", "별도 3D 제작 대상일 때만 분리\n쿠션·가드·헤드레스트", PURPLE, PURPLE_LIGHT),
    ]
    for idx, (code, title, body, color, light) in enumerate(lanes):
        x = 0.7 + idx * 2.43
        add_box(slide, x, 3.34, 2.2, 1.55, fill=light, line=color)
        add_pill(slide, code, x + 0.17, 3.56, 0.62, fill=color, color=WHITE, size=8)
        add_text(slide, title, x + 0.17, 4.02, 1.82, 0.24, size=10.1, color=color, bold=True)
        add_text(slide, body, x + 0.17, 4.37, 1.86, 0.4, size=7.8, color=INK)

    add_box(slide, 0.7, 5.22, 7.55, 1.33, fill=BLUE_LIGHT, line=BLUE)
    add_text(slide, "확정 사례 · G25070005743", 0.98, 5.46, 2.25, 0.24, size=10.5, color=BLUE_DARK, bold=True)
    add_text(slide, "seq 1  4인 소파", 1.0, 5.88, 1.35, 0.22, size=9.2, bold=True)
    add_text(slide, "2910 / 1020 / 910 mm", 2.32, 5.88, 1.9, 0.22, size=9.2, color=BLUE_DARK, bold=True)
    add_text(slide, "seq 2  스툴", 4.42, 5.88, 1.15, 0.22, size=9.2, bold=True)
    add_text(slide, "740 / 660 / 410 mm", 5.62, 5.88, 1.85, 0.22, size=9.2, color=BLUE_DARK, bold=True)
    add_text(slide, "같은 product_id, 다른 component_seq", 1.0, 6.2, 5.95, 0.2, size=8.2, color=MUTED)

    add_box(slide, 8.52, 5.22, 4.11, 1.33, fill=GREEN_LIGHT, line=GREEN)
    add_text(slide, "침대 예외 정책", 8.82, 5.46, 1.5, 0.24, size=10.5, color=GREEN, bold=True)
    add_text(
        slide,
        f"프레임+매트리스는 단일 Asset {stats.bed_single_asset:,}개\n단, SS/Q/K처럼 외형 크기가 달라지면 option_seq 분리",
        8.82,
        5.84,
        3.25,
        0.48,
        size=8.7,
        color=INK,
        bold=True,
    )


def slide_pattern_queue(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Pattern work queue",
        f"체크필요 {stats.pattern_check_needed:,}개는 두 패턴 체계로 쪼개 순차 승인",
        page,
        "D 코드는 규격 후보 구조, N/S/Q/M/A/O는 3D Asset 구조를 뜻합니다. D00은 조합 구조를 먼저 판단합니다.",
    )

    add_box(slide, 0.7, 1.78, 5.82, 4.25, fill=WHITE, line=AMBER)
    add_pill(slide, f"규격 패턴 · {stats.comparison:,}개", 0.98, 2.03, 1.75, fill=AMBER, color=WHITE, size=9)
    add_text(slide, "RAW 후보를 대표·옵션·제외값으로 분류", 2.9, 2.08, 3.1, 0.22, size=9.2, color=MUTED)
    dim_defs = [
        ("D00", "규격+조합 동시 확인", stats.dimension_patterns.get("D00", 0)),
        ("D01", "완전 W/D/H 1세트", stats.dimension_patterns.get("D01", 0)),
        ("D02", "완전 1세트+부분 후보", stats.dimension_patterns.get("D02", 0)),
        ("D03", "완전 W/D/H 복수", stats.dimension_patterns.get("D03", 0)),
        ("D04", "부분 후보 1세트", stats.dimension_patterns.get("D04", 0)),
        ("D05", "부분 후보 복수", stats.dimension_patterns.get("D05", 0)),
    ]
    for idx, (code, label, count) in enumerate(dim_defs):
        yy = 2.62 + idx * 0.49
        add_pill(slide, code, 1.0, yy, 0.58, fill=AMBER_LIGHT, color=AMBER, size=8)
        add_text(slide, label, 1.76, yy + 0.04, 2.65, 0.22, size=9.1, bold=True)
        add_text(slide, f"{count:,}", 4.58, yy + 0.01, 0.8, 0.25, size=11, color=AMBER, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, "개", 5.44, yy + 0.05, 0.28, 0.2, size=8.4, color=MUTED)

    add_box(slide, 6.78, 1.78, 5.85, 4.25, fill=WHITE, line=LG_RED)
    add_pill(
        slide,
        f"조합 패턴 · {stats.pattern_combination + stats.pattern_both:,}개",
        7.06,
        2.03,
        1.62,
        fill=LG_RED,
        color=WHITE,
        size=9,
    )
    add_text(slide, "한 상품에 여러 패턴군이 동시에 존재 가능", 8.85, 2.08, 3.1, 0.22, size=9.2, color=MUTED)
    combo_defs = [
        ("N", "조합 탐지 제외", stats.combination_pattern_groups.get("N", 0)),
        ("S", "완성된 단일 조립체", stats.combination_pattern_groups.get("S", 0)),
        ("Q", "동일 Asset 수량 묶음", stats.combination_pattern_groups.get("Q", 0)),
        ("M", "서로 다른 Asset 분리", stats.combination_pattern_groups.get("M", 0)),
        ("A", "부속품 별도 제작 판단", stats.combination_pattern_groups.get("A", 0)),
        ("O", "형상·규격 옵션 분기", stats.combination_pattern_groups.get("O", 0)),
    ]
    for idx, (code, label, count) in enumerate(combo_defs):
        yy = 2.62 + idx * 0.49
        add_pill(slide, code, 7.08, yy, 0.58, fill=LG_PINK, color=LG_RED, size=8)
        add_text(slide, label, 7.84, yy + 0.04, 2.65, 0.22, size=9.1, bold=True)
        add_text(slide, f"{count:,}", 10.66, yy + 0.01, 0.8, 0.25, size=11, color=LG_RED, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, "개", 11.52, yy + 0.05, 0.28, 0.2, size=8.4, color=MUTED)

    add_box(slide, 0.7, 6.2, 11.93, 0.58, fill=INK, line=INK)
    add_text(
        slide,
        "승인 흐름  대표 페이지 확인 → 패턴 규칙·예외 확정 → rule_version 저장 → ledger 재산출 → Excel 상태 갱신 → 회귀 테스트",
        1.02,
        6.37,
        11.25,
        0.22,
        size=8.7,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_req1_lineage(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Request 1 · field lineage",
        "요청 1의 모든 셀은 ‘원천 → 정규화 → 값·상태·근거’ 순서로 산출",
        page,
        "개발자는 값만 저장하지 말고 어느 원천과 규칙으로 만들어졌는지 함께 추적할 수 있어야 합니다.",
    )
    headers = ["요청 필드", "우선 원천", "정규화·출력", "상태 판정"]
    rows = [
        [
            "대표·롤링 이미지",
            "goods API image list\nsortSeq=1,2,3…",
            "대표 URL + rolling_url_1~18\nURL은 원문 그대로 저장",
            "URL 존재=확보\n각도·배경 품질은 별도 검수",
        ],
        [
            "W / D / H (mm)",
            "API 고시 → HTML/FAQ/Q&A\n→ 상세 이미지 OCR",
            "cm×10 · inch×25.4\nW/D/H 숫자형 + 원문/후보",
            "확정값 / RAW분리 / 무후보\n패턴 상태는 별도 컬럼",
        ],
        [
            "분류·추천 공간",
            "goods category\nspace collection API",
            "mid/small category\n리빙룸|베드룸… 표준어",
            "API 명시=확보\n카테고리 규칙=추론 표시",
        ],
        [
            "브랜드·색상·옵션",
            "brandName · 고시정보\npurchaseOptions/stock",
            "brand · color\noption_style + option_1…N",
            "명시값=확보\n추론값=노란색, 없음=분홍색",
        ],
        [
            "설치·배치 위치",
            "가전 설치 속성\n카테고리·명시 문구",
            "빌트인/스탠딩\n벽·천장·바닥 + 추천높이",
            "홈스타일 설치 타입=해당없음\n추천높이 미명시=미확보",
        ],
    ]
    row_fills = [BLUE_LIGHT, AMBER_LIGHT, WHITE, WHITE, GRAY_LIGHT]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.72,
        11.93,
        4.08,
        col_widths=[1.9, 3.0, 3.35, 3.68],
        header_fill=BLUE_DARK,
        font_size=8.45,
        first_col_bold=True,
        row_fills=row_fills,
    )
    add_box(slide, 0.7, 6.02, 11.93, 0.72, fill=BLUE_LIGHT, line=rgb("#C7D9FA"))
    add_text(slide, "권장 DB 저장 단위", 0.98, 6.22, 1.45, 0.22, size=9.5, color=BLUE_DARK, bold=True)
    add_text(
        slide,
        "product_id · field_code · value · value_type · unit · status · source_type · source_key · evidence · rule_version · updated_at",
        2.55,
        6.19,
        9.55,
        0.26,
        size=8.5,
        color=INK,
        bold=True,
    )


def slide_req2_lineage(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Request 2 · implementation levels",
        "요청 2는 하나의 AI 결과가 아니라 네 가지 산출 수준으로 구현",
        page,
        "현재는 LLM·embedding을 사용하지 않습니다. 명시값과 규칙값, 관계 후보, 미수집을 반드시 구분합니다.",
    )
    levels = [
        ("1", "명시 추출", "API·HTML·FAQ/Q&A에\n실제로 존재하는 값", BLUE, BLUE_LIGHT),
        ("2", "규칙 추론", "상품명·설명·분류를\n사전 키워드로 태깅", AMBER, AMBER_LIGHT),
        ("3", "관계 후보", "같은 공간에 등장한\n제품·스타일 연결", GREEN, GREEN_LIGHT),
        ("4", "미수집", "CRM·공간 크기처럼\n현재 원천에 없는 값", MUTED, GRAY_LIGHT),
    ]
    for idx, (num, title, body, color, light) in enumerate(levels):
        x = 0.7 + idx * 3.04
        add_box(slide, x, 1.7, 2.8, 1.1, fill=light, line=color)
        add_pill(slide, num, x + 0.2, 1.94, 0.42, fill=color, color=WHITE, size=8)
        add_text(slide, title, x + 0.78, 1.93, 1.65, 0.24, size=11, color=color, bold=True)
        add_text(slide, body, x + 0.2, 2.32, 2.3, 0.34, size=7.9, color=INK)

    headers = ["고객 요청", "현재 제공값 / 수준", "개발자가 알아야 할 경계"]
    rows = [
        [
            "설명서 13개 영역 태깅",
            "13개 컬럼 · 명시 추출",
            "상품고시 우선, FAQ/Q&A 보강\n빈 문자열을 추론값으로 채우지 않음",
        ],
        [
            "디자인 스타일 추론",
            "미니멀·곡선형·모듈형 등 · 규칙 추론",
            "이미지 판정이나 LLM 결과가 아님\n스타일 사전·rule_version 필요",
        ],
        [
            "공간 콘텐츠 태깅",
            "스타일·분위기·공간·색상톤 · 명시+추론",
            "space API와 카테고리 규칙을 구분\n공간 크기는 현재 미수집",
        ],
        [
            "공간스타일 ↔ 제품",
            "space/style ↔ product_id · 관계 후보",
            "현재는 collection 동시 등장 관계\nedge table·score·evidence로 구조화 필요",
        ],
        [
            "제품 ↔ 제품",
            "product_id ↔ product_id · 관계 후보",
            "동시 등장은 배치 이유나 추천 점수가 아님\nrelation_type·방향·근거 필요",
        ],
        [
            "개인 구매·CRM",
            "미수집",
            "공개 PDP 범위에 개인 데이터 없음\n권한·동의 후 별도 CRM 연계",
        ],
        [
            "의미기반 검색·추천",
            "검색용 문장 · 전처리 결과",
            "레드|적색|붉은색 등 동의어 결합\nvector index·검색 API는 아직 미구현",
        ],
    ]
    row_fills = [WHITE, AMBER_LIGHT, GREEN_LIGHT, BLUE_LIGHT, BLUE_LIGHT, GRAY_LIGHT, AMBER_LIGHT]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        3.04,
        11.93,
        3.34,
        col_widths=[2.25, 3.45, 6.23],
        header_fill=GREEN,
        font_size=8.05,
        first_col_bold=True,
        row_fills=row_fills,
    )
    add_box(slide, 0.7, 6.53, 11.93, 0.3, fill=GREEN_LIGHT, line=rgb("#B7DEBF"))
    add_text(
        slide,
        "저장 계약: product_id · field_code · value · status(EXPLICIT/RULE/RELATION_CANDIDATE/NOT_COLLECTED) · source · evidence · confidence",
        1.0,
        6.59,
        11.2,
        0.16,
        size=7.6,
        color=GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_quality(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Quality gate",
        "현재 결과는 DB·파서·Excel을 서로 대조해 재현성과 정합성을 확인",
        page,
        "테스트 통과는 현재 수집 완료 DB 기준이며, 빈 DB 통합 재구축 테스트는 아직 추가해야 합니다.",
    )
    checks = [
        ("SQLite integrity", "ok", "DB 파일 구조 정상"),
        ("Foreign key", "0 오류", "고아 참조 없음"),
        ("상품 원장", f"{stats.total:,} / {stats.total:,}", "누락·중복 0"),
        ("파서·피드백 회귀", "PASS", "정규화·오탐 사례"),
        ("Excel 구조", f"{stats.workbook_sheets} / {stats.workbook_columns}", "시트 / 주시트 열"),
        ("패턴 상태 행", f"{stats.total:,}", "필터·상태 합계 일치"),
        ("비교 후보", f"{stats.comparison_candidates:,}", "DB와 Excel 일치"),
        ("스타일·정렬", "0 오류", "노란셀·후보 위치"),
    ]
    positions = [(0.7 + (i % 4) * 3.04, 1.93 + (i // 4) * 1.55) for i in range(8)]
    for (label, value, note), (x, y) in zip(checks, positions):
        add_box(slide, x, y, 2.8, 1.25, fill=WHITE, line=LINE)
        add_pill(slide, "PASS", x + 0.2, y + 0.18, 0.6, fill=GREEN_LIGHT, color=GREEN, size=7.7)
        add_text(slide, label, x + 0.95, y + 0.19, 1.58, 0.25, size=9.5, bold=True)
        add_text(slide, value, x + 0.2, y + 0.61, 1.3, 0.34, size=18, color=GREEN, bold=True)
        add_text(slide, note, x + 1.22, y + 0.7, 1.3, 0.2, size=7.6, color=MUTED, align=PP_ALIGN.RIGHT)

    add_box(slide, 0.7, 5.2, 11.93, 1.05, fill=AMBER_LIGHT, line=rgb("#F2C36B"))
    add_text(slide, "현재 테스트 범위 밖", 1.0, 5.49, 1.65, 0.25, size=11, color=AMBER, bold=True)
    add_text(
        slide,
        "빈 폴더 bootstrap · API schema fixture · OCR 이미지 fixture · 수집→OCR→원장→Excel 통합 테스트",
        2.72,
        5.45,
        9.25,
        0.31,
        size=10.5,
        color=INK,
    )


def slide_readiness(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Code readiness",
        "PoC 배치로는 사용 가능, 운영 배치로 전환하려면 P0 두 항목을 먼저 해소",
        page,
        "현재 DB를 입력으로 한 후처리·Excel 재생성은 가능하지만, 새 PC·빈 DB에서 end-to-end 한 번 실행은 보장하지 않습니다.",
    )
    headers = ["영역", "상태", "현재 판단", "운영 전 조치"]
    rows = [
        ["수집 API/HTML/Q&A", "완료", "9,358개 원천 저장", "schema 감시·SSL·retry 지표"],
        ["OCR·규격 원장", "완료", "2단계 OCR·단일 상태 원장", "run version·resume"],
        ["Excel 재생성/검증", "완료", f"{stats.workbook_sheets}시트·{stats.workbook_columns}열·상태필터", "CI 자동 검증"],
        ["빈 DB 최초 구축", "미완성", "순환 의존성", "bootstrap mode"],
        ["단일 실행기·설정", "미완성", "단계별 수동 실행", "pipeline.yaml + runner"],
        ["설치 명세", "미완성", "환경에 의존", "requirements·OCR preflight"],
        ["scope 변경 동기화", "부분", "upsert만 수행", "active snapshot/validity"],
        ["운영 서비스/API", "미구현", "SQLite+Excel 배치", "field/evidence API"],
    ]
    fills = [
        GREEN_LIGHT,
        GREEN_LIGHT,
        GREEN_LIGHT,
        RED_LIGHT,
        RED_LIGHT,
        AMBER_LIGHT,
        AMBER_LIGHT,
        GRAY_LIGHT,
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.92,
        11.93,
        4.75,
        col_widths=[2.25, 1.25, 3.5, 4.93],
        header_fill=INK,
        font_size=9.15,
        first_col_bold=True,
        row_fills=fills,
    )


def slide_roadmap(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Productization roadmap",
        "P0에서 패턴·상태를 고정하고, P1에서 재구축·운영, P2에서 추천 활용을 확장",
        page,
        "일정은 개발 리소스와 운영 방식 결정 후 산정하며, 아래 순서는 의존성 기준입니다.",
    )
    roadmap = [
        (
            "P0",
            "규칙·상태 고정",
            "D00~D05 승인\nN/S/Q/M/A/O 승인\nrule_version·회귀 테스트",
            LG_RED,
            LG_PINK,
        ),
        (
            "P1",
            "재구축·운영",
            "bootstrap 순환 해소\n단일 runner·config\nrun ID·resume·scope snapshot",
            BLUE,
            BLUE_LIGHT,
        ),
        (
            "P2",
            "서비스 연계",
            "최종 field/evidence fact\nAPI·로그·모니터링\nembedding·관계 edge",
            GREEN,
            GREEN_LIGHT,
        ),
    ]
    x_values = [0.7, 4.55, 8.4]
    for idx, ((priority, title, body, color, light), x) in enumerate(zip(roadmap, x_values)):
        add_box(slide, x, 2.0, 3.55, 3.55, fill=WHITE, line=color)
        add_pill(slide, priority, x + 0.28, 2.28, 0.62, fill=color, color=WHITE, size=9)
        add_text(slide, title, x + 1.1, 2.29, 1.85, 0.32, size=16, color=color, bold=True)
        parts = body.split("\n")
        add_bullet_list(
            slide,
            parts,
            x + 0.32,
            3.14,
            2.87,
            1.55,
            size=11.2,
            bullet_color=color,
            gap=0.54,
        )
        gate_label = ["Gate: 체크필요 상태 전환", "Gate: 빈 DB → Excel", "Gate: API·검색 소비"][idx]
        add_box(slide, x + 0.28, 4.88, 2.99, 0.39, fill=light, line=light)
        add_text(slide, gate_label, x + 0.37, 4.98, 2.8, 0.18, size=8.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        if idx < 2:
            add_arrow(slide, x + 3.6, 3.75, x_values[idx + 1] - 0.06, 3.75, MUTED)

    add_box(slide, 0.7, 5.91, 11.93, 0.7, fill=INK, line=INK)
    add_text(slide, "권장 착수 기준", 1.0, 6.13, 1.25, 0.22, size=10, color=WHITE, bold=True)
    add_text(
        slide,
        "P0 완료 전에는 현재 DB를 기준으로만 재산출하고, 네트워크 수집·OCR 재실행 전 DB/OCR 폴더를 버전 백업",
        2.42,
        6.1,
        9.45,
        0.27,
        size=9.7,
        color=rgb("#E1E5EA"),
    )


def slide_execution_overview(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix",
        "현재 DB 재산출은 ‘후보 보강 → 선택 → 원장 → 판정 → Excel 검증’의 5단계",
        page,
        "이 경로는 네트워크와 이미지 OCR을 다시 호출하지 않고, 이미 저장된 DB 후보를 결정론적으로 재평가합니다.",
    )
    preflight = [
        ("DB 경로", "bulk_homestyle_collect.DB_PATH", "현재 homestyle_bulk.sqlite"),
        ("입력 상태", "products / sources / OCR staging", "수집·OCR 결과가 이미 존재"),
        ("실행 방식", "SQLite WAL + run/snapshot", "기존 raw 원천은 변경하지 않음"),
        ("사전 조치", "DB 백업 · Excel 종료", "파일 잠금과 재실행 혼선 방지"),
    ]
    for idx, (title, code, note) in enumerate(preflight):
        x = 0.7 + idx * 3.04
        add_box(slide, x, 1.9, 2.8, 1.12, fill=WHITE, line=LINE)
        add_text(slide, title, x + 0.2, 2.12, 0.82, 0.22, size=9, color=LG_RED, bold=True)
        add_text(slide, code, x + 0.2, 2.43, 2.35, 0.22, size=8.5, bold=True)
        add_text(slide, note, x + 0.2, 2.71, 2.35, 0.18, size=7.5, color=MUTED)

    stages = [
        ("01", "파서 회귀", "test_dimension_\ncontext_normalizer.py", LG_RED),
        ("02", "비교 후보 보강", "partial / blocked /\npass2 raw exposure", BLUE),
        ("03", "현재 선택 재구축", "refresh_dimension_\ncontext_selection.py", TEAL),
        ("04", "단일 원장 갱신", "build_dimension_\nresolution_ledger.py ×2", AMBER),
        ("05", "패턴·Excel·검증", "remaining → pattern\n→ component → workbook", GREEN),
    ]
    x_values = [0.7, 3.19, 5.68, 8.17, 10.66]
    for idx, (num, title, code, color) in enumerate(stages):
        add_box(slide, x_values[idx], 3.55, 1.98, 1.65, fill=WHITE, line=color)
        add_pill(slide, num, x_values[idx] + 0.18, 3.77, 0.43, fill=color, color=WHITE, size=7.5)
        add_text(slide, title, x_values[idx] + 0.18, 4.17, 1.58, 0.26, size=10.1, bold=True)
        add_text(slide, code, x_values[idx] + 0.18, 4.57, 1.57, 0.43, size=7.8, color=MUTED)
        if idx < 4:
            add_arrow(slide, x_values[idx] + 2.01, 4.38, x_values[idx + 1] - 0.04, 4.38, LINE)

    add_box(slide, 0.7, 5.58, 7.63, 1.07, fill=BLUE_LIGHT, line=rgb("#C7D9FA"))
    add_text(slide, "재실행 안전장치", 1.0, 5.84, 1.35, 0.23, size=10.5, color=BLUE_DARK, bold=True)
    add_text(
        slide,
        "후보 생성 파일은 동일 run_name 행을 삭제 후 재삽입하고,\n판정 테이블은 새 snapshot을 만든 뒤 이전 is_current를 0으로 전환합니다.",
        2.42,
        5.78,
        5.45,
        0.5,
        size=9.2,
        color=INK,
    )
    add_box(slide, 8.58, 5.58, 4.05, 1.07, fill=RED_LIGHT, line=rgb("#F0B4B2"))
    add_text(slide, "이 경로에서 하지 않는 일", 8.88, 5.8, 3.1, 0.23, size=10.5, color=RED, bold=True)
    add_text(slide, "API 재수집 · 이미지 재다운로드 · Windows OCR 재실행 · scope 변경", 8.88, 6.13, 3.25, 0.34, size=8.2, color=INK)


def slide_execution_candidates(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix · Phase A",
        "후보 보강 파일은 값을 확정하지 않고 ‘사람이 비교할 정보’를 추가",
        page,
        "세 파일 모두 stg_dimension_targeted_ocr_candidate에 run_name별 결과를 다시 만들며, 자동 잠금은 이후 원장이 결정합니다.",
    )
    headers = ["파일", "왜 실행하는가", "주요 입력", "출력·정상 확인"]
    rows = [
        [
            "test_dimension_context_normalizer.py",
            "규격 정규식·축 매핑·OCR 보정 회귀를 먼저 차단",
            "코드 내 21개 testcase\n(DB 사용 안 함)",
            "21개 PASS\n하나라도 실패하면 이후 단계 중단",
        ],
        [
            "run_dimension_partial_fusion.py",
            "같은 이미지의 W/H, D/H 등 부분 축을 W/D/H 조합으로 결합",
            "classification master·ledger\nremaining targeted 후보",
            "run=dimension_partial_fusion_v1\nfusion_candidates / truncated_groups",
        ],
        [
            "run_dimension_blocked_candidate_exposure.py",
            "완성 W/D/H가 있지만 자동 선택이 막힌 값을 비교 후보로 노출",
            "미확정 ledger\n현재 context complete 후보",
            "run=dimension_blocked_candidate_exposure_v1\nexposed_candidates",
        ],
        [
            "run_dimension_pass2_raw_exposure.py",
            "layout OCR의 축 미확정 숫자 묶음도 버리지 않고 비교용으로 노출",
            "pass2 product_image\n+ pass2 observation",
            "run=dimension_pass2_raw_exposure_v1\nraw / unit_inferred / truncated",
        ],
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.92,
        11.93,
        3.98,
        col_widths=[2.95, 3.0, 2.85, 3.13],
        header_fill=BLUE_DARK,
        font_size=8.35,
        first_col_bold=True,
        row_fills=[LG_PINK, BLUE_LIGHT, BLUE_LIGHT, BLUE_LIGHT],
    )
    add_box(slide, 0.7, 6.08, 11.93, 0.65, fill=AMBER_LIGHT, line=rgb("#F2C36B"))
    add_text(
        slide,
        "공통 정책: decision_status=HUMAN_REVIEW · 비교 후보는 자동 PASS 규격으로 사용 금지 · 2D 제품은 잘못된 W/D/H 합성을 건너뜀",
        1.0,
        6.27,
        11.15,
        0.23,
        size=9.2,
        color=AMBER,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_execution_context_ledger(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix · Phase B",
        "refresh는 현재 후보 선택을 재구축하고, ledger는 최종 상태를 누적·잠금",
        page,
        "두 파일의 책임을 섞지 않는 것이 중요합니다. refresh는 staging 계산, ledger는 authoritative fact와 이력을 관리합니다.",
    )
    add_box(slide, 0.7, 1.9, 5.78, 4.8, fill=WHITE, line=TEAL)
    add_pill(slide, "STEP 05", 1.0, 2.16, 0.78, fill=TEAL, color=WHITE, size=8)
    add_text(slide, "refresh_dimension_context_selection.py", 1.0, 2.64, 4.65, 0.3, size=13.5, bold=True, color=TEAL)
    add_text(slide, "현재 context snapshot 안에서 후보 선택을 다시 계산", 1.0, 3.03, 4.7, 0.26, size=9.4, color=MUTED)
    add_bullet_list(
        slide,
        [
            "현재 stg_dimension_context_product의 snapshot_id 탐색",
            "targeted 5개 run 후보를 context candidate로 병합",
            "상품별 대표 후보·복수 옵션·REOCR queue 재선택",
            "회귀 testcase 결과를 같은 snapshot에 다시 기록",
        ],
        1.02,
        3.48,
        4.92,
        1.65,
        size=10.2,
        bullet_color=TEAL,
        gap=0.41,
    )
    add_box(slide, 1.0, 5.37, 5.05, 0.92, fill=TEAL_LIGHT, line=TEAL_LIGHT)
    add_text(slide, "WRITE", 1.18, 5.56, 0.55, 0.2, size=8, color=TEAL, bold=True)
    add_text(
        slide,
        "stg_dimension_context_candidate / product\nstg_product_dimension_option / targeted_reocr_queue / regression",
        1.85,
        5.48,
        3.9,
        0.48,
        size=8.3,
        color=INK,
    )

    add_box(slide, 6.85, 1.9, 5.78, 4.8, fill=WHITE, line=AMBER)
    add_pill(slide, "STEP 06", 7.15, 2.16, 0.78, fill=AMBER, color=WHITE, size=8)
    add_text(slide, "build_dimension_resolution_ledger.py", 7.15, 2.64, 4.7, 0.3, size=13.5, bold=True, color=AMBER)
    add_text(slide, "source 확정·rule 확정·비교·무후보를 상품별 한 행으로 고정", 7.15, 3.03, 4.8, 0.26, size=9.4, color=MUTED)
    add_bullet_list(
        slide,
        [
            "classification master + context + option + 기존 ledger 결합",
            "MANUAL/기존 비교/유효한 잠금은 진단 snapshot이 내려도 보존",
            "상태·대표값이 바뀔 때만 hist event 기록",
            "확정 옵션과 비교 후보를 별도 fact에 보존",
        ],
        7.17,
        3.48,
        4.92,
        1.65,
        size=10.2,
        bullet_color=AMBER,
        gap=0.41,
    )
    add_box(slide, 7.15, 5.37, 5.05, 0.92, fill=AMBER_LIGHT, line=AMBER_LIGHT)
    add_text(slide, "WRITE", 7.33, 5.56, 0.55, 0.2, size=8, color=AMBER, bold=True)
    add_text(
        slide,
        "fact_dimension_resolution_ledger / option / comparison_candidate\nhist_dimension_resolution_event",
        8.0,
        5.48,
        3.92,
        0.48,
        size=8.3,
        color=INK,
    )
    add_box(slide, 4.72, 6.43, 3.88, 0.35, fill=RED_LIGHT, line=rgb("#F0B4B2"))
    add_text(slide, "ledger는 반드시 2회 실행 → 2회차 transitions_written = 0", 4.9, 6.51, 3.5, 0.18, size=8.3, color=RED, bold=True, align=PP_ALIGN.CENTER)


def slide_execution_outputs(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix · Phase C",
        "원장이 안정된 뒤 잔여·패턴·구성품·필수값·Excel을 순서대로 생성",
        page,
        "뒤 단계는 앞 단계의 현재 snapshot을 대조하므로, 파일만 따로 실행하면 stale snapshot 오류가 발생할 수 있습니다.",
    )
    headers = ["순서 / 파일", "처리 방식", "DB·파일 출력", "정상 종료 기준"]
    rows = [
        [
            "07  build_dimension_remaining_\nreason_staging.py",
            "ledger의 NO_CANDIDATE만 대상으로 pass1/pass2/crop/candidate 존재 여부를 6개 원인으로 분류",
            "stg_dimension_remaining_reason\n이전 current=0, 새 snapshot=1",
            f"rows={stats.no_candidate:,}\nreason 합계={stats.no_candidate:,}\nintegrity_check=ok",
        ],
        [
            "08  combination pattern\n+ bed asset policy",
            "조합 미확정 149개 패턴과 침대 단일 Asset/사이즈 옵션 정책을 current snapshot으로 생성",
            "stg_combination_candidate_pattern\nstg_bed_asset_size_option",
            f"조합 체크={stats.pattern_combination + stats.pattern_both:,}\n침대 단일 Asset={stats.bed_single_asset:,}\n패턴군 코드 생성",
        ],
        [
            "09  product component\n+ mandatory staging",
            "조합상품 구성품 행을 만든 뒤 별첨 6개 필수 그룹의 PASS snapshot을 DB에 저장",
            "stg_product_component_dimension\nstg_mandatory_pass",
            f"상품={stats.total:,}\n구성품 rows={stats.component_rows:,}\nPASS={stats.mandatory_pass:,}",
        ],
        [
            "10  build_homestyle_\nbulk_workbook.py",
            "현재 mandatory·조합상품 snapshot과 메모리 계산이 같을 때만 9개 시트 작성",
            "…요구필드_대량결과_\n패턴상태.xlsx",
            f"sheets={stats.workbook_sheets}\nrows={stats.total:,} / columns={stats.workbook_columns}\n구성품 rows={stats.component_rows:,}",
        ],
        [
            "11  validate_dimension_\ncomparison_workbook.py",
            "XLSX XML을 직접 읽어 규격 열 인접성, 쉼표 후보 정렬, DB 후보 수·work queue를 대조",
            "파일 변경 없음\n검증 JSON만 stdout",
            f"candidate_sum={stats.comparison_candidates:,}\nalignment_errors=0\nDB integrity=ok",
        ],
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.92,
        11.93,
        4.72,
        col_widths=[2.75, 4.15, 2.65, 2.38],
        header_fill=GREEN,
        font_size=7.7,
        first_col_bold=True,
        row_fills=[WHITE, AMBER_LIGHT, GREEN_LIGHT, GREEN_LIGHT, WHITE],
    )
    add_text(
        slide,
        "중요: mandatory staging은 workbook의 build_rows()에 의존하므로, 빈 DB bootstrap 시 발생하는 순환 의존성의 핵심 지점입니다.",
        0.85,
        6.73,
        11.55,
        0.2,
        size=7.9,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_execution_checklist(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix · Runbook",
        "복사해 실행할 명령과 단계별 중단 기준",
        page,
        "PowerShell에서 PYTHONIOENCODING을 UTF-8로 지정하고 프로젝트 루트에서 순서대로 실행합니다.",
    )
    commands = [
        ("01", "python test_dimension_context_normalizer.py"),
        ("02", "python run_dimension_partial_fusion.py"),
        ("03", "python run_dimension_blocked_candidate_exposure.py"),
        ("04", "python run_dimension_pass2_raw_exposure.py"),
        ("05", "python refresh_dimension_context_selection.py"),
        ("06A", "python build_dimension_resolution_ledger.py"),
        ("06B", "python build_dimension_resolution_ledger.py  # idempotency 확인"),
        ("07", "python build_dimension_remaining_reason_staging.py"),
        ("08", "python build_combination_candidate_pattern_staging.py"),
        ("09", "python build_bed_asset_size_option_staging.py"),
        ("10", "python build_product_component_staging.py"),
        ("11", "python build_mandatory_pass_staging.py"),
        ("12", "python build_homestyle_bulk_workbook.py"),
        ("13", "python validate_dimension_comparison_workbook.py"),
    ]
    for idx, (num, cmd) in enumerate(commands):
        col = 0 if idx < 7 else 1
        row = idx if idx < 7 else idx - 7
        x = 0.7 + col * 6.08
        y = 1.82 + row * 0.45
        add_box(slide, x, y, 5.85, 0.34, fill=WHITE, line=LINE)
        add_pill(slide, num, x + 0.1, y, 0.46, fill=INK, color=WHITE, size=6.5)
        add_text(slide, cmd, x + 0.68, y + 0.07, 4.92, 0.2, size=7.7, bold=True)

    checks = [
        (
            "즉시 중단",
            "test FAIL · current context snapshot 없음 · SQLite integrity ≠ ok",
            RED,
            RED_LIGHT,
        ),
        (
            "원장 확인",
            f"2회차 transition=0 · 상태 합계={stats.total:,} · locked={stats.locked:,} · comparison={stats.comparison:,}",
            AMBER,
            AMBER_LIGHT,
        ),
        (
            "최종 확인",
            f"패턴 체크={stats.pattern_check_needed:,} · workbook={stats.workbook_sheets} sheets/{stats.workbook_columns} cols\n후보={stats.comparison_candidates:,} · 정렬 오류=0",
            GREEN,
            GREEN_LIGHT,
        ),
    ]
    for idx, (title, body, color, light) in enumerate(checks):
        x = 0.7 + idx * 4.06
        add_box(slide, x, 5.1, 3.82, 1.15, fill=light, line=color)
        add_text(slide, title, x + 0.25, 5.29, 2.9, 0.24, size=10.2, color=color, bold=True)
        add_text(slide, body, x + 0.25, 5.68, 3.3, 0.42, size=8.0, color=INK)
    add_box(slide, 0.7, 6.47, 11.93, 0.32, fill=INK, line=INK)
    add_text(
        slide,
        "$env:PYTHONIOENCODING='utf-8'  ·  실행 전 SQLite와 homestyle_bulk_run/ocr 폴더를 실행 버전명으로 백업",
        1.0,
        6.55,
        11.2,
        0.18,
        size=8,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_code_change_guide(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Developer appendix · Change guide",
        "요구사항 변경 시 수정 위치를 계층별로 제한해야 원장과 Excel이 어긋나지 않음",
        page,
        "fact 테이블을 직접 수정하지 말고, 규칙·후보·snapshot을 변경한 뒤 동일 runbook으로 원장을 다시 만듭니다.",
    )
    headers = ["변경 목적", "주 수정 파일", "반드시 함께 확인할 것"]
    rows = [
        [
            "OCR 오인식·표기 정규식 추가",
            "dimension_context_normalizer.py",
            "test_dimension_context_normalizer.py에 재현 testcase 추가",
        ],
        [
            "대표 후보·복수 옵션 선택 변경",
            "build_dimension_context_normalization.py\nrefresh_dimension_context_selection.py",
            "17개 사용자 회귀 · current context/option/queue 수",
        ],
        [
            "비교 후보 노출 정책 변경",
            "run_dimension_*_exposure.py\nrun_dimension_partial_fusion.py",
            "HUMAN_REVIEW 유지 · candidate 상한·중복·단위 추론",
        ],
        [
            "잠금·상태 전이·무효화 변경",
            "build_dimension_resolution_ledger.py",
            "hist event · 기존 MANUAL/LOCK 보존 · 2회차 transition=0",
        ],
        [
            "D00~D05 규격 패턴 분류 변경",
            "build_homestyle_bulk_workbook.py\n(dimension_review_pattern_code)",
            "3,411 합계 · D코드 상호배타 · D00 조합 후보 교집합",
        ],
        [
            "N/S/Q/M/A/O 조합 패턴 변경",
            "build_combination_candidate_pattern_staging.py\nbuild_homestyle_bulk_workbook.py",
            "149개 고유 상품 · 다중 패턴 허용 · 옵션 O군 중복 제거",
        ],
        [
            "필수 PASS 기준 변경",
            "build_homestyle_bulk_workbook.py\nbuild_mandatory_pass_staging.py",
            "STANDARD_VERSION · 대상수/충족수 · summary/missing view",
        ],
        [
            "요청 1·2 출력 필드 추가",
            "build_homestyle_bulk_workbook.py\nvalidate_dimension_comparison_workbook.py",
            "데이터사전 · 열 색상 · 숫자/단위 · DB/Excel 정렬 검증",
        ],
    ]
    add_table(
        slide,
        headers,
        rows,
        0.7,
        1.92,
        11.93,
        4.42,
        col_widths=[2.7, 4.1, 5.13],
        header_fill=PURPLE,
        font_size=7.7,
        first_col_bold=True,
        row_fills=[PURPLE_LIGHT, WHITE, WHITE, AMBER_LIGHT, AMBER_LIGHT, LG_PINK, GREEN_LIGHT, BLUE_LIGHT],
    )
    add_box(slide, 0.7, 6.47, 11.93, 0.32, fill=RED_LIGHT, line=rgb("#F0B4B2"))
    add_text(
        slide,
        "금지: fact_dimension_resolution_ledger 직접 UPDATE · 기존 run_name 재사용 상태에서 raw 후보 혼합 · validator 생략 후 Excel 배포",
        0.95,
        6.55,
        11.4,
        0.18,
        size=8,
        color=RED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def slide_decision(prs: Presentation, stats: Stats, page: int):
    slide = add_slide_base(
        prs,
        "Decision",
        "다음 목표: 패턴 체크 3,484개를 승인 규칙으로 전환하고 재현 가능한 배치로 운영",
        page,
        "234개 무후보 보강과 3,484개 패턴 확정은 다른 작업이며, 상태 전이와 회귀 검증을 분리 운영합니다.",
    )
    add_box(slide, 0.7, 1.95, 11.93, 1.05, fill=INK, line=INK)
    add_text(
        slide,
        f"{stats.total:,}개 산출 상태 100% 분류 · 확정값 {stats.locked:,}개 · 패턴 체크 {stats.pattern_check_needed:,}개\n별첨 필수값 PASS {stats.mandatory_pass:,}개({stats.pass_pct:.1f}%)",
        1.03,
        2.17,
        8.6,
        0.57,
        size=15,
        color=WHITE,
        bold=True,
    )
    add_pill(slide, "POC VERIFIED", 10.24, 2.29, 1.58, fill=LG_RED, color=WHITE, size=9)

    decisions = [
        ("1", "패턴 승인 회의", "D00→D05 및 N/S/Q/M/A/O 순서로 규칙 확정"),
        ("2", "상태 전이 구현", "RAW→확정/무후보 · 조합→seq 행 · rule_version"),
        ("3", "P0 제품화", "bootstrap·단일 runner·통합 smoke test"),
        ("4", "잔여·가전 분리", "무후보 234 보강 / 가전 ontology 별도"),
    ]
    for idx, (num, title, note) in enumerate(decisions):
        x = 0.7 + (idx % 2) * 6.08
        y = 3.42 + (idx // 2) * 1.45
        add_box(slide, x, y, 5.85, 1.13, fill=WHITE, line=LINE)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.24), Inches(y + 0.29), Inches(0.5), Inches(0.5)
        )
        set_fill(circle, LG_RED)
        circle.line.fill.background()
        add_text(
            slide,
            num,
            x + 0.24,
            y + 0.32,
            0.5,
            0.32,
            size=11,
            color=WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_text(slide, title, x + 0.97, y + 0.23, 2.6, 0.28, size=13, bold=True)
        add_text(slide, note, x + 0.97, y + 0.62, 4.2, 0.24, size=9.2, color=MUTED)

    add_text(
        slide,
        "산출 파일: 홈스타일_비음영대상군_전체상품_요구필드_대량결과_패턴상태.xlsx",
        0.72,
        6.63,
        7.3,
        0.2,
        size=8,
        color=MUTED,
    )
    add_text(
        slide,
        "기준 DB: homestyle_bulk_run/homestyle_bulk.sqlite",
        8.05,
        6.63,
        4.3,
        0.2,
        size=8,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def build_deck() -> Path:
    if not DB_PATH.exists():
        raise FileNotFoundError(f"DB not found: {DB_PATH}")
    stats = load_stats()
    if stats.total != 9358:
        raise RuntimeError(f"Unexpected product count: {stats.total}")
    if stats.source_confirmed + stats.rule_resolved + stats.comparison + stats.no_candidate != stats.total:
        raise RuntimeError("Dimension status total does not match product total")
    if stats.mandatory_pass + stats.reinforcement != stats.total:
        raise RuntimeError("Mandatory status total does not match product total")
    if stats.pattern_no_check + stats.pattern_check_needed != stats.total:
        raise RuntimeError("Pattern status total does not match product total")
    if sum(stats.dimension_patterns.values()) != stats.comparison:
        raise RuntimeError("Dimension pattern total does not match RAW total")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "홈스타일 상품정보 수집·규격 OCR 개발 인수인계 보고"
    prs.core_properties.subject = "API·HTML·FAQ/Q&A·OCR 기반 요청 1·2 필드 산출"
    prs.core_properties.author = "HomeStyle Data PoC"
    prs.core_properties.keywords = "HomeStyle, OCR, API, product data, developer handoff"
    prs.core_properties.comments = "DB와 workbook 산출 메타를 재조회하여 2026-07-27 기준으로 생성"

    make_cover(prs, stats)
    slide_executive(prs, stats, 2)
    slide_requests(prs, stats, 3)
    slide_input_output(prs, stats, 4)
    slide_sources(prs, stats, 5)
    slide_ocr(prs, stats, 6)
    slide_rules(prs, stats, 7)
    slide_dimension_result(prs, stats, 8)
    slide_mandatory(prs, stats, 9)
    slide_remaining(prs, stats, 10)
    slide_examples(prs, stats, 11)
    slide_workbook(prs, stats, 12)
    slide_excel_status_model(prs, stats, 13)
    slide_3d_row_policy(prs, stats, 14)
    slide_pattern_queue(prs, stats, 15)
    slide_req1_lineage(prs, stats, 16)
    slide_req2_lineage(prs, stats, 17)
    slide_quality(prs, stats, 18)
    slide_readiness(prs, stats, 19)
    slide_roadmap(prs, stats, 20)
    slide_execution_overview(prs, stats, 21)
    slide_execution_candidates(prs, stats, 22)
    slide_execution_context_ledger(prs, stats, 23)
    slide_execution_outputs(prs, stats, 24)
    slide_execution_checklist(prs, stats, 25)
    slide_code_change_guide(prs, stats, 26)
    slide_decision(prs, stats, 27)

    prs.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_deck()
    print(path)
